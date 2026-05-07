"""Reproduce two ChatGPT-style layouts as editable PowerPoint slides.

Slide 1 — `fig_pipeline_v3.pptx`:
    7-stage construction pipeline (Source data → Search region sampling →
    Target & evidence selection → Evidence alignment → Shortcut-aware
    filtering → Distractor generation → Final benchmark item) with a
    coloured numbered badge per stage.

Slide 2 — `fig_teaser_v4.pptx`:
    Refined teaser following the second ChatGPT image, but with the
    Charlie / Charlie St. Cloud example (good_examples/example_01).

Both decks are built so every panel, badge, arrow, and text block is its
own pptx shape — the user can click and edit anything inside PowerPoint
or Keynote, and re-running this script overwrites both files in place.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path("/jumbo/jinlab/Sally/ForSeeBench")
FE = ROOT / "figure_examples"
OUT = ROOT / "figures/paper_figures_pptx"
OUT.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------------------
# Palette (slightly extended to cover seven stage colours)
# ---------------------------------------------------------------------------
def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


WHITE = _rgb("#FFFFFF")
INK = _rgb("#1F2933")
MUTED = _rgb("#6B7280")
SUBTLE = _rgb("#9CA3AF")

# Per-stage colours (matching the ChatGPT pipeline image)
STAGE_COLORS = [
    ("#2D6FB6", "#E8F0FB"),   # 1 Source Data — blue
    ("#D97A2D", "#FCEEDB"),   # 2 Search Region Sampling — orange
    ("#3F8F5C", "#DDEFD2"),   # 3 Target & Evidence Selection — green
    ("#1B7F8E", "#DCEEEC"),   # 4 Evidence Alignment — teal
    ("#C0413E", "#FBE3E1"),   # 5 Shortcut-Aware Filtering — red
    ("#E2A83C", "#FCEEDB"),   # 6 Distractor Generation — gold
    ("#7E5BB0", "#EFE7F4"),   # 7 Final Benchmark Item — purple
]


# Charlie example data (same as fig_teaser_v3)
EX_FOLDER = FE / "good_examples/example_01"
AD_LINES = [
    ("AD₁", "“Now, Charlie strides out of the restaurant, his eyes cast downward.”"),
    ("AD₂", "“As Charlie steps into the street, a fisherman in waterproof bib pants passes him carrying a coil of rope slung over his shoulder.”"),
]
BRIDGE_INDEX = 1
TARGET_TEXT = "“Charlie stares mesmerized out at the water.”"
QUESTION = "Q:  What happens next?"
OPTIONS = [
    ("A", "Charlie gazes intently at the water.",                  True),
    ("B", "Charlie looks around the street for a lost pet.",       False),
    ("C", "Charlie picks up a piece of trash on the sidewalk.",    False),
    ("D", "Charlie talks to a passerby about the weather.",        False),
]


def _pick_frames(folder: Path, n: int = 3) -> list[Path]:
    out = []
    for tag in ["t0", "tmid", "tend"]:
        for fp in sorted(folder.glob("frames/context_*.jpg")):
            if fp.name.endswith(f"__{tag}.jpg"):
                out.append(fp)
                break
    return out[:n]


def _pick_target_frame(folder: Path) -> Path | None:
    matches = sorted(folder.glob("frames/target__*tmid*.jpg"))
    return matches[0] if matches else None


# ---------------------------------------------------------------------------
# Shape helpers (kept compatible with python-pptx)
# ---------------------------------------------------------------------------
def add_rect(slide, x, y, w, h, *, fill, line, line_w=1.0, rounded=True,
             corner=0.10):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    if rounded:
        s.adjustments[0] = corner
    if fill is not None:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_circle(slide, cx, cy, r, *, fill, line=None, line_w=1.0):
    s = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r),
        Inches(2 * r), Inches(2 * r),
    )
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, color, size, bold=False,
             italic=False, align="left", vertical="middle", monospace=False,
             line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if vertical == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif vertical == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = "Consolas" if monospace else "DejaVu Sans"
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color
    return box


def add_arrow(slide, x1, y1, x2, y2, *, color, width=2.0, dashed=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    line.line.color.rgb = color; line.line.width = Pt(width)
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = line.line._get_or_add_ln()
    tail_end = ln.find(qn("a:tailEnd"))
    if tail_end is None:
        tail_end = etree.SubElement(ln, qn("a:tailEnd"))
    tail_end.set("type", "triangle")
    tail_end.set("w", "med"); tail_end.set("len", "med")
    if dashed:
        prst_dash = ln.find(qn("a:prstDash"))
        if prst_dash is None:
            prst_dash = etree.SubElement(ln, qn("a:prstDash"))
        prst_dash.set("val", "dash")
    return line


# ---------------------------------------------------------------------------
# Slide 1 — Construction pipeline (7 stages, ChatGPT image 1)
# ---------------------------------------------------------------------------
def build_pipeline() -> None:
    prs = Presentation()
    W, H = 16.0, 9.0   # widescreen for a horizontal pipeline
    prs.slide_width = Inches(W); prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title — matches §4 "Benchmark Construction" of the paper.
    add_text(slide, 0.0, 0.18, W, 0.45,
             "ForSeeBench: Benchmark Construction",
             color=INK, size=24, bold=True, align="center")
    add_text(slide, 0.0, 0.62, W, 0.30,
             "10 MAD-eval movies  →  6,520 ordered AD segments  →  787 prospective QA items",
             color=MUTED, size=12, italic=True, align="center")

    # ---- Stage layout (3 rows × 3 cols + 1 wide bottom row) ----
    # Six top-row stages share a 3x2 grid; stage 7 is full width below.
    margin = 0.30
    top_rows = 2; top_cols = 3
    bottom_h = 1.85
    grid_top = 0.95
    grid_bot = H - bottom_h - 0.30 - 0.40    # leave room for legend

    grid_h = grid_bot - grid_top - 0.30
    cell_w = (W - 2 * margin - 0.20 * (top_cols - 1)) / top_cols
    cell_h = (grid_h - 0.20 * (top_rows - 1)) / top_rows

    # Paper-faithful stages: paragraph headers are taken verbatim from
    # `paper/.../sec/dataset.tex`, and the symbols c_i, c_y, C_y, q_y match
    # `paper/.../sec/task_formulation.tex` (`Let c_1, ..., c_n denote an
    # ordered sequence of AD segments...`).
    # Order matches the paper's "Construction pipeline" paragraph: question
    # generation (stage 4 in the paper) precedes validation (stage 5).
    stage_specs = [
        # 1.  Source data  (paper paragraph: "Source data")
        {
            "title": "Source data",
            "lines": [
                "10 MAD-eval movies.",
                "Per movie, ordered AD segments",
                "    c₁ , c₂ , … , cₙ",
                "(human-written; each cᵢ has a timestamp).",
            ],
            "subnote": "Total: 6,520 AD segments across 10 movies.",
        },
        # 2.  Candidate search and target selection  (paper paragraph)
        {
            "title": "Candidate search and target selection",
            "lines": [
                "Partition the ordered AD segments into",
                "10-clip search regions",
                "    {cᵢ , cᵢ₊₁ , … , cᵢ₊₉}.",
                "After a valid target is selected the search",
                "advances past it; otherwise it advances by 1 clip.",
            ],
            "subnote": "A search region is a construction-time window,\nnot the context shown to the answering model.",
        },
        # 3.  Target selection and evidence  (paper paragraph)
        {
            "title": "Target selection and evidence  (Qwen)",
            "lines": [
                "Qwen2.5-VL-7B-Instruct picks within each region:",
                "    • a hidden target  cᵧ",
                "    • a supporting prior context  Cᵧ ⊂ {c₁ , … , cᵧ₋₁}",
                "    • exact evidence spans inside Cᵧ.",
                "Assigns target type and evidence-sufficiency score.",
            ],
            "subnote": "target_type ∈ {action_transition, state_change,\nparticipant_update, spatial_consequence,\nobject_reveal, visible_text_update}.",
        },
        # 4.  Evidence alignment  (Construction pipeline step 3 in the paper)
        {
            "title": "Evidence alignment",
            "lines": [
                "Reconcile every selected id in Cᵧ and every",
                "evidence span against the original AD stream.",
                "Reject any span that does not match strictly",
                "prior AD text (i.e., index i < y).",
            ],
            "subnote": "Guarantees temporal leakage is excluded by\nconstruction (no clip with i ≥ y enters Cᵧ).",
        },
        # 5.  Questions and distractors  (paper paragraph)
        {
            "title": "Questions and distractors",
            "lines": [
                "From (cᵧ, Cᵧ, evidence) Qwen writes one question",
                "    qᵧ : “What happens next?”",
                "and four options: 1 correct + 3 typed distractors.",
                "Distractor labels:",
                "    already_happened · entity_swapped",
                "    plausible_unsupported · contradicts_context · unrelated",
            ],
            "subnote": "Typed distractors make answerer errors interpretable.",
        },
        # 6.  Validation and audit metadata  (paper paragraph)
        {
            "title": "Validation and audit metadata",
            "lines": [
                "Automated checks: target answerability, context support,",
                "distractor plausibility, no-context / movie-prior /",
                "dialogue-only / wrong-context leakage.",
                "Retain only items with",
                "    Qwen confidence ≥ 0.7,",
                "    evidence sufficiency ≥ 0.7,",
                "    distractor quality ≥ 0.7.",
            ],
            "subnote": "Every rejection is stored with an inspectable reason\n(insufficient evidence · triviality · weak distractors · …).",
        },
    ]

    badge_color = lambda i: _rgb(STAGE_COLORS[i][0])
    soft_color = lambda i: _rgb(STAGE_COLORS[i][1])

    # Compute cell positions
    cell_xs = [margin + c * (cell_w + 0.20) for c in range(top_cols)]
    cell_ys = [grid_top + r * (cell_h + 0.20) for r in range(top_rows)]

    for i, spec in enumerate(stage_specs):
        r = i // top_cols; c = i % top_cols
        x = cell_xs[c]; y = cell_ys[r]
        ec = badge_color(i); fc = soft_color(i)

        # Outer cell — soft tint with coloured border
        add_rect(slide, x, y, cell_w, cell_h,
                 fill=WHITE, line=ec, line_w=1.4, corner=0.06)

        # Numbered circle badge in the top-left
        badge_r = 0.22
        add_circle(slide, x + 0.30, y + 0.28, badge_r, fill=ec)
        add_text(slide, x + 0.30 - badge_r, y + 0.28 - badge_r,
                 2 * badge_r, 2 * badge_r,
                 str(i + 1), color=WHITE, size=14, bold=True,
                 align="center", vertical="middle")

        # Stage title (right of the badge)
        add_text(slide, x + 0.65, y + 0.05, cell_w - 0.80, 0.50,
                 spec["title"],
                 color=ec, size=13, bold=True, vertical="middle")

        # Body lines
        body_text = "\n".join(spec["lines"])
        add_text(slide, x + 0.25, y + 0.65, cell_w - 0.50, cell_h - 1.15,
                 body_text,
                 color=INK, size=10.5, vertical="top",
                 line_spacing=1.30)

        # Subnote ribbon at the bottom of the cell
        sub_h = 0.45
        add_rect(slide, x + 0.20, y + cell_h - sub_h - 0.05,
                 cell_w - 0.40, sub_h,
                 fill=fc, line=None, corner=0.20)
        add_text(slide, x + 0.30, y + cell_h - sub_h - 0.05,
                 cell_w - 0.60, sub_h,
                 spec["subnote"],
                 color=ec, size=9.5, italic=True, vertical="middle",
                 line_spacing=1.20)

    # Arrows between cells in reading order (1→2→3 across, then 4→5→6 across).
    # Use thin grey arrows so the cells are the focus.
    for i in range(len(stage_specs) - 1):
        r1, c1 = i // top_cols, i % top_cols
        r2, c2 = (i + 1) // top_cols, (i + 1) % top_cols
        if r1 == r2:
            x1 = cell_xs[c1] + cell_w
            x2 = cell_xs[c2]
            y_mid = cell_ys[r1] + cell_h / 2
            add_arrow(slide, x1 - 0.05, y_mid, x2 + 0.05, y_mid,
                      color=SUBTLE, width=2.0)
        else:
            # Wrap from end of row r1 to start of row r2.
            x1 = cell_xs[c1] + cell_w / 2
            y1 = cell_ys[r1] + cell_h + 0.05
            x2 = cell_xs[c2] + cell_w / 2
            y2 = cell_ys[r2] - 0.05
            # A simple zig-zag using two arrows.
            add_arrow(slide, x1, y1, x1, (y1 + y2) / 2,
                      color=SUBTLE, width=2.0)
            add_arrow(slide, x1, (y1 + y2) / 2, x2, y2,
                      color=SUBTLE, width=2.0)

    # ---- Stage 7 — bottom-row "Final Benchmark Item" (full width) ----
    f_x = margin
    f_y = grid_bot
    f_w = W - 2 * margin
    ec = badge_color(6); fc = soft_color(6)
    add_rect(slide, f_x, f_y, f_w, bottom_h,
             fill=WHITE, line=ec, line_w=1.6, corner=0.04)
    # Badge
    add_circle(slide, f_x + 0.30, f_y + 0.30, 0.22, fill=ec)
    add_text(slide, f_x + 0.30 - 0.22, f_y + 0.08, 0.44, 0.44,
             "7", color=WHITE, size=14, bold=True, align="center",
             vertical="middle")
    add_text(slide, f_x + 0.65, f_y + 0.05, f_w - 0.80, 0.45,
             "Final benchmark item   ·   Release format and statistics",
             color=ec, size=14, bold=True, vertical="middle")

    # 3 columns of content inside the final-item card
    inner_top = f_y + 0.55
    inner_h = bottom_h - 0.65
    col_w = (f_w - 0.50 * 4) / 3
    col_xs = [f_x + 0.30 + i * (col_w + 0.50) for i in range(3)]

    # Column A — Visible to the answerer
    cx = col_xs[0]
    add_text(slide, cx, inner_top, col_w, 0.30,
             "Visible to the answerer",
             color=ec, size=11, bold=True, vertical="top")
    add_text(slide, cx, inner_top + 0.32, col_w, inner_h - 0.32,
             "Prior context  Cᵧ ⊂ {c₁ , … , cᵧ₋₁}\n"
             "Question  qᵧ  : “What happens next?”\n"
             "Four options  (A) (B) (C) (D)",
             color=INK, size=10.5, vertical="top", line_spacing=1.40)

    # Column B — Hidden / used only for scoring
    cx = col_xs[1]
    add_text(slide, cx, inner_top, col_w, 0.30,
             "Hidden / used only for scoring",
             color=ec, size=11, bold=True, vertical="top")
    add_text(slide, cx, inner_top + 0.32, col_w, inner_h - 0.32,
             "Target AD sentence  cᵧ\n"
             "Correct option = paraphrase of cᵧ\n"
             "Evidence spans inside Cᵧ supporting cᵧ",
             color=INK, size=10.5, vertical="top", line_spacing=1.40)

    # Column C — Stored audit metadata
    cx = col_xs[2]
    add_text(slide, cx, inner_top, col_w, 0.30,
             "Stored audit metadata",
             color=ec, size=11, bold=True, vertical="top")
    add_text(slide, cx, inner_top + 0.32, col_w, inner_h - 0.32,
             "•  target type  (1 of 6)\n"
             "•  distractor labels  (1 correct + 3 of 5 typed labels)\n"
             "•  reasoning type, continuity type\n"
             "•  Qwen confidence / evidence sufficiency / distractor quality\n"
             "•  rejection reasons for filtered candidates",
             color=INK, size=10.5, vertical="top", line_spacing=1.40)

    # ---- Legend (bottom row) ----
    legend_y = H - 0.40
    add_text(slide, margin, legend_y, 1.7, 0.30,
             "Legend (stage colours):",
             color=INK, size=10, bold=True, vertical="middle")
    legend_items = [
        ("1 Source data",                     STAGE_COLORS[0][0]),
        ("2 Candidate search",                STAGE_COLORS[1][0]),
        ("3 Target + evidence",               STAGE_COLORS[2][0]),
        ("4 Evidence alignment",              STAGE_COLORS[3][0]),
        ("5 Question + distractors",          STAGE_COLORS[4][0]),
        ("6 Validation + audit",              STAGE_COLORS[5][0]),
        ("7 Final item",                      STAGE_COLORS[6][0]),
    ]
    lx = margin + 1.8
    for label, hex_col in legend_items:
        add_circle(slide, lx + 0.10, legend_y + 0.15, 0.10,
                   fill=_rgb(hex_col))
        add_text(slide, lx + 0.30, legend_y, 1.6, 0.30,
                 label,
                 color=INK, size=9.5, vertical="middle")
        lx += 1.85

    out = OUT / "fig_pipeline_v3.pptx"
    prs.save(out)
    print(f"  wrote {out}")


# ---------------------------------------------------------------------------
# Slide 2 — Refined teaser (ChatGPT image 2) using the Charlie example
# ---------------------------------------------------------------------------
def build_teaser_v4() -> None:
    prs = Presentation()
    W, H = 13.5, 8.5
    prs.slide_width = Inches(W); prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    TITLE_BAR = _rgb("#1B5E5E")
    PANEL_GREY = _rgb("#F1F1F2")
    AM_TEAL = _rgb("#2D7D8C"); AM_HI = _rgb("#D9EAEE")
    AT_OR = _rgb("#E78B3A")
    PRIOR_HI = _rgb("#FFF3D2")
    KW_GREEN = _rgb("#1F8B5E")
    GREEN = _rgb("#3F8F5C"); GREEN_HI = _rgb("#DDEFD2")
    DARK_GREEN = _rgb("#1E5631")
    BLUE = _rgb("#2D6FB6")
    PINK = _rgb("#E29093")
    RED = _rgb("#C0413E"); RED_HI = _rgb("#FBE3E1")
    DARK_RED = _rgb("#7B1F1C")

    # 1) Title bar
    title_h = 0.55
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(W), Inches(title_h),
    )
    bg.fill.solid(); bg.fill.fore_color.rgb = TITLE_BAR
    bg.line.fill.background(); bg.shadow.inherit = False
    add_text(slide, 0, 0, W, title_h,
             "ForSeeBench: Prospective Audio-Description QA",
             color=WHITE, size=22, bold=True, align="center", vertical="middle")

    # 2) AD-generation pipeline strip
    pipe_y = title_h + 0.10
    pipe_h = 1.85
    add_rect(slide, 0.20, pipe_y, W - 0.40, pipe_h,
             fill=PANEL_GREY, line=SUBTLE, line_w=0.8, rounded=False)
    add_text(slide, 0.35, pipe_y + 0.10, 6.0, 0.30,
             "Existing AD-generation pipeline",
             color=INK, size=12, bold=True, vertical="top")

    # frames
    fr_y = pipe_y + 0.45
    fr_h = pipe_h - 0.65
    fr_w = 1.40
    fr_gap = 0.10
    fr_x_start = 0.40
    frames = _pick_frames(EX_FOLDER, 3)
    for i, fp in enumerate(frames):
        slide.shapes.add_picture(str(fp),
                                 Inches(fr_x_start + i * (fr_w + fr_gap)),
                                 Inches(fr_y),
                                 Inches(fr_w), Inches(fr_h))

    a1_x = fr_x_start + 3 * (fr_w + fr_gap) - fr_gap
    arr_y = fr_y + fr_h / 2
    add_arrow(slide, a1_x + 0.10, arr_y, a1_x + 0.55, arr_y,
              color=AM_TEAL, width=2.4)

    # AD Model box
    am_x = a1_x + 0.65
    am_w = 1.85
    add_rect(slide, am_x, fr_y, am_w, fr_h,
             fill=AM_HI, line=AM_TEAL, line_w=1.6, corner=0.10)
    add_text(slide, am_x, fr_y + 0.05, am_w, 0.50, "AD Model",
             color=AM_TEAL, size=14, bold=True, align="center", vertical="middle")
    add_text(slide, am_x, fr_y + 0.55, am_w, fr_h - 0.65,
             "(MAD-eval, NarrAD,\nAutoAD-Zero)",
             color=AM_TEAL, size=10, align="center", vertical="top",
             line_spacing=1.30)

    add_arrow(slide, am_x + am_w + 0.10, arr_y, am_x + am_w + 0.55, arr_y,
              color=AM_TEAL, width=2.4)

    # AD text bubble (orange-bordered)
    at_x = am_x + am_w + 0.65
    at_w = W - 0.20 - at_x - 0.10
    add_rect(slide, at_x, fr_y, at_w, fr_h,
             fill=WHITE, line=AT_OR, line_w=1.6, corner=0.10)
    add_text(slide, at_x + 0.15, fr_y + 0.05, at_w - 0.30, 0.30,
             "AD text  (generated)",
             color=AT_OR, size=12, bold=True, vertical="top")
    yy = fr_y + 0.45
    for tag, line in AD_LINES:
        add_text(slide, at_x + 0.15, yy, at_w - 0.30, 0.45,
                 line, color=INK, size=10, vertical="top",
                 line_spacing=1.20)
        yy += 0.45

    # 3) Firewall divider
    fw_y = pipe_y + pipe_h + 0.20
    fw_line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.20), Inches(fw_y), Inches(W - 0.20), Inches(fw_y),
    )
    fw_line.line.color.rgb = GREEN
    fw_line.line.width = Pt(1.8)
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = fw_line.line._get_or_add_ln()
    pd = ln.find(qn("a:prstDash"))
    if pd is None:
        pd = etree.SubElement(ln, qn("a:prstDash"))
    pd.set("val", "dash")
    add_text(slide, 0, fw_y + 0.05, W, 0.32,
             "▼  The answerer reads only what is below this line.  ▼",
             color=GREEN, size=14, bold=True, italic=True, align="center")

    # 4) Body — two columns (smaller bottom info bar reserved)
    body_y_top = fw_y + 0.55
    bottom_bar_h = 0.85
    body_y_bot = H - bottom_bar_h - 0.30
    body_h = body_y_bot - body_y_top

    left_x = 0.20; left_w = 7.50
    right_x = left_x + left_w + 0.20
    right_w = W - right_x - 0.20

    # ── LEFT: Prior AD context (visible to the model) ──
    add_rect(slide, left_x, body_y_top, left_w, body_h,
             fill=WHITE, line=BLUE, line_w=1.6, corner=0.05)
    # Header strip
    add_rect(slide, left_x, body_y_top, left_w, 0.45,
             fill=BLUE, line=None, corner=0.20)
    add_text(slide, left_x + 0.20, body_y_top, left_w - 0.40, 0.45,
             "Prior AD context  (visible to the model)",
             color=WHITE, size=13, bold=True, vertical="middle")

    line_y = body_y_top + 0.65
    line_h = 0.95
    for i, (tag, text) in enumerate(AD_LINES):
        is_bridge = (i == BRIDGE_INDEX)
        if is_bridge:
            add_rect(slide, left_x + 0.20, line_y - 0.05,
                     left_w - 0.40, line_h + 0.10,
                     fill=PRIOR_HI, line=None, corner=0.06)
        add_text(slide, left_x + 0.30, line_y, 0.80, 0.40, tag,
                 color=AM_TEAL, size=14, bold=True, monospace=True,
                 vertical="top")
        add_text(slide, left_x + 1.10, line_y, left_w - 1.30, line_h,
                 text, color=INK, size=11.5, vertical="top",
                 line_spacing=1.30)
        line_y += line_h + 0.10

    # Bridge keywords callout
    callout_y = line_y - 0.05
    add_text(slide, left_x + 0.30, callout_y, 2.20, 0.30,
             "▸ bridge keywords:",
             color=DARK_RED, size=11, bold=True, italic=True, vertical="top")
    add_text(slide, left_x + 2.55, callout_y, left_w - 2.85, 0.30,
             "fisherman   ·   waterproof   ·   rope",
             color=KW_GREEN, size=12, bold=True, vertical="top")

    # Bottom note inside the left column
    add_text(slide, left_x + 0.30, body_y_bot - 0.42, left_w - 0.60, 0.30,
             "Infer the next narrated visual update from this AD line alone.",
             color=MUTED, size=11, italic=True, vertical="middle")

    # ── RIGHT TOP: Correct next AD (hidden) ──
    rt_h = body_h * 0.42
    rt_y_top = body_y_top
    add_rect(slide, right_x, rt_y_top, right_w, rt_h,
             fill=RED_HI, line=PINK, line_w=1.6, corner=0.06)
    # Header strip
    add_rect(slide, right_x, rt_y_top, right_w, 0.40,
             fill=PINK, line=None, corner=0.20)
    add_text(slide, right_x, rt_y_top, right_w, 0.40,
             "Correct next AD  (hidden)",
             color=WHITE, size=12, bold=True, align="center", vertical="middle")
    # 🔒 lock icon on the right of the header strip
    add_text(slide, right_x + right_w - 0.50, rt_y_top, 0.45, 0.40,
             "🔒",
             color=WHITE, size=14, bold=True, align="center", vertical="middle")

    target_fp = _pick_target_frame(EX_FOLDER)
    fr_x_t = right_x + 0.20
    fr_y_t = rt_y_top + 0.55
    fr_w_t = 1.35
    fr_h_t = rt_h - 0.75
    if target_fp:
        slide.shapes.add_picture(str(target_fp),
                                 Inches(fr_x_t), Inches(fr_y_t),
                                 Inches(fr_w_t), Inches(fr_h_t))
    tx = fr_x_t + fr_w_t + 0.18
    tw = right_x + right_w - tx - 0.20
    add_text(slide, tx, fr_y_t, tw, fr_h_t,
             TARGET_TEXT,
             color=DARK_RED, size=11, italic=True, vertical="top",
             line_spacing=1.30)

    # ── RIGHT BOTTOM: Question + options ──
    rb_y_top = rt_y_top + rt_h + 0.20
    rb_h = body_y_bot - rb_y_top
    add_rect(slide, right_x, rb_y_top, right_w, rb_h,
             fill=WHITE, line=GREEN, line_w=1.6, corner=0.06)
    add_text(slide, right_x + 0.20, rb_y_top + 0.05, right_w - 0.40, 0.40,
             QUESTION,
             color=INK, size=13, bold=True, vertical="middle")

    opt_top = rb_y_top + 0.55
    opt_h = (rb_h - 0.65) / 4
    for j, (letter, text, is_correct) in enumerate(OPTIONS):
        oy = opt_top + j * opt_h
        if is_correct:
            add_rect(slide, right_x + 0.15, oy, right_w - 0.30, opt_h - 0.05,
                     fill=GREEN_HI, line=None, corner=0.10)
        add_text(slide, right_x + 0.25, oy, 0.50, opt_h - 0.05,
                 f"({letter})",
                 color=GREEN if is_correct else INK, size=11, bold=True,
                 vertical="middle")
        add_text(slide, right_x + 0.75, oy, right_w - 1.20, opt_h - 0.05,
                 text, color=INK, size=10, vertical="middle")
        if is_correct:
            add_text(slide, right_x + right_w - 0.45, oy,
                     0.30, opt_h - 0.05,
                     "✓",
                     color=GREEN, size=14, bold=True,
                     align="right", vertical="middle")

    # 5) Bottom info bar — 4 categorical chips
    bar_y = body_y_bot + 0.18
    bar_h = bottom_bar_h - 0.36
    chips = [
        ("Visible to model",      "prior AD context only",                  BLUE),
        ("Hidden from model",     "next AD (future event)",                 PINK),
        ("Model output",          "multiple-choice MCQ\ncorrect next AD",   AM_TEAL),
        ("Goal",                  "predict context evidence\nfor next narrated visual update", GREEN),
        ("No shortcuts",          "crafted distractors prevent\ntrivial solutions", _rgb("#7E5BB0")),
    ]
    n_chips = len(chips)
    chip_gap = 0.18
    chip_w = (W - 0.60 - chip_gap * (n_chips - 1)) / n_chips
    chip_x = 0.30
    for label, body, color in chips:
        # left bullet circle
        add_circle(slide, chip_x + 0.18, bar_y + bar_h * 0.30, 0.10, fill=color)
        add_text(slide, chip_x + 0.35, bar_y, chip_w - 0.40, bar_h * 0.45,
                 label,
                 color=color, size=10.5, bold=True, vertical="middle")
        add_text(slide, chip_x + 0.35, bar_y + bar_h * 0.45,
                 chip_w - 0.40, bar_h * 0.55,
                 body,
                 color=INK, size=9, vertical="top", line_spacing=1.20)
        chip_x += chip_w + chip_gap

    out = OUT / "fig_teaser_v4.pptx"
    prs.save(out)
    print(f"  wrote {out}")


def main() -> None:
    print("Building ChatGPT-style PPTX figures…")
    build_pipeline()
    build_teaser_v4()
    print("Done.")


if __name__ == "__main__":
    main()
