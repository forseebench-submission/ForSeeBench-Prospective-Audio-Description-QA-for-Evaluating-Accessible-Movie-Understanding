"""Build a clean three-row schematic for the ForSeeBench context conditions.

Outputs:
  figures/paper_figures/fig_context_conditions.{pdf,png}
  figures/paper_figures_pptx/fig_context_conditions.pptx

Three rows on a shared clip timeline:
  (i)  No context  (k=0):   no prior clip reaches the answerer.
  (ii) Fixed window (k=4):  the last k clips reach the answerer.
  (iii) Adaptive context:   only evidence-selected clips reach the answerer
                            (non-contiguous; not necessarily the tail).

The target clip is rendered in red on every row to make visible that it is
hidden from the answerer in all three settings.
"""
from __future__ import annotations

from pathlib import Path

import matplotlib.patches as patches
import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch

ROOT = Path("/jumbo/jinlab/Sally/ForSeeBench")
OUT_PNG = ROOT / "figures/paper_figures"
OUT_PPTX = ROOT / "figures/paper_figures_pptx"
OUT_PNG.mkdir(parents=True, exist_ok=True)
OUT_PPTX.mkdir(parents=True, exist_ok=True)

INK       = "#1F2933"
MUTED     = "#6B7280"
SUBTLE    = "#9CA3AF"
LIGHT     = "#E5E7EB"
BLUE      = "#3D6FB6"
BLUE_HI   = "#D9E5F4"
ORANGE    = "#D97A2D"
ORANGE_HI = "#FCE5C9"
GREEN     = "#3F8F5C"
RED       = "#C0413E"
DARK_RED  = "#7B1F1C"
RED_HI    = "#FBE3E1"


def _clip(ax, x, y, w, h, *, fc, ec, lw=1.0, label=None,
          label_color=INK, label_size=8.0, label_bold=False, alpha=1.0,
          hatch=None):
    """Draw one clip cell, optionally with a hatch fill (used for the hidden
    target box to visually signal 'withheld')."""
    ax.add_patch(FancyBboxPatch(
        (x, y), w, h,
        boxstyle="round,pad=0.0,rounding_size=0.012",
        lw=lw, ec=ec, fc=fc, alpha=alpha, zorder=2, hatch=hatch,
    ))
    if label:
        ax.text(x + w / 2, y + h / 2, label,
                ha="center", va="center", fontsize=label_size,
                color=label_color, fontweight="bold" if label_bold else "normal",
                zorder=3)


def make_matplotlib() -> None:
    plt.rcParams.update({
        "font.family": "DejaVu Sans",
        "font.size": 9,
        "pdf.fonttype": 42, "ps.fonttype": 42,
    })

    # Wider canvas + more breathing room. Final figure is ~5.4 × 3.4".
    fig = plt.figure(figsize=(5.4, 3.5))
    fig.patch.set_facecolor("white")
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")

    # Title
    ax.text(0.50, 0.965,
            "Context conditions",
            ha="center", va="center", fontsize=11, fontweight="bold", color=INK)
    ax.text(0.50, 0.928,
            "Same item, same hidden target — only which prior AD clips reach the answerer changes.",
            ha="center", va="center", fontsize=7.4, color=MUTED, style="italic")

    # Geometry. Reserve a wide left gutter for row labels so they never collide
    # with the clip cells.
    label_x = 0.030       # left edge of row label text
    label_w = 0.250       # space reserved for label + sublabel
    timeline_left = 0.300 # first clip starts here
    timeline_right = 0.770
    target_x = 0.815
    target_w = 0.155
    n_clips = 6
    clip_gap = 0.010
    cell_w = (timeline_right - timeline_left - clip_gap * (n_clips - 1)) / n_clips
    row_h = 0.085

    rows = [
        (0.745,
         "(i)  No context",
         "k = 0",
         set(),
         MUTED, LIGHT,
         "Answerer receives no prior AD.",
         "→ 0 clips"),
        (0.480,
         "(ii)  Fixed window",
         "k = 4",
         {2, 3, 4, 5},
         BLUE, BLUE_HI,
         "Answerer receives the most recent k AD clips.",
         "→ 4 clips"),
        (0.215,
         "(iii)  Adaptive",
         "evidence-selected",
         {1, 4},
         ORANGE, ORANGE_HI,
         "Answerer receives only the AD clips selected as evidence at\nconstruction time — may skip clips, not always the tail.",
         "→ 2 clips"),
    ]

    for (y, label, sublabel, selected, accent, accent_hi,
         caption, tag) in rows:
        # Row label
        ax.text(label_x, y + row_h - 0.012, label,
                fontsize=9.4, fontweight="bold", color=INK, va="top")
        ax.text(label_x, y + row_h - 0.052, sublabel,
                fontsize=7.4, color=accent if selected else MUTED,
                va="top", style="italic", fontweight="bold")

        # Clip cells
        for i in range(n_clips):
            x = timeline_left + i * (cell_w + clip_gap)
            is_selected = i in selected
            if is_selected:
                _clip(ax, x, y, cell_w, row_h,
                      fc=accent_hi, ec=accent, lw=1.6,
                      label=f"AD{_sub(i+1)}",
                      label_color=accent, label_bold=True, label_size=8.4)
            else:
                _clip(ax, x, y, cell_w, row_h,
                      fc="#F5F5F7", ec=SUBTLE, lw=0.8, alpha=0.85,
                      label=f"AD{_sub(i+1)}",
                      label_color=SUBTLE, label_size=7.6)

        # Target cell — hatched red box, no "?" symbol. The hatch + the
        # explicit "(hidden)" label are how we visually communicate that the
        # next AD is withheld in every condition.
        _clip(ax, target_x, y, target_w, row_h,
              fc=RED_HI, ec=RED, lw=1.4, hatch="///",
              label=None)
        ax.text(target_x + target_w / 2, y + row_h / 2 + 0.012,
                "next AD",
                ha="center", va="center", fontsize=7.6,
                color=DARK_RED, fontweight="bold", zorder=3)
        ax.text(target_x + target_w / 2, y + row_h / 2 - 0.018,
                "(hidden)",
                ha="center", va="center", fontsize=6.4,
                color=DARK_RED, style="italic", zorder=3)

        # Caption (italic, INK) — left-aligned under the clip strip.
        ax.text(timeline_left + 0.005, y - 0.018, caption,
                fontsize=6.8, color=INK, va="top", style="italic",
                linespacing=1.30)

        # Tag (accent, bold) — right-aligned under the target column, on the
        # same baseline as the caption so they don't overlap horizontally.
        ax.text(0.985, y - 0.018, tag,
                fontsize=6.8, color=accent, ha="right", va="top",
                style="italic", fontweight="bold")

    # Footer
    ax.text(0.50, 0.040,
            "The next AD sentence is hidden in every condition; only the prior subset given to the answerer changes.",
            ha="center", va="center", fontsize=7.0, color=MUTED, style="italic")

    pdf = OUT_PNG / "fig_context_conditions.pdf"
    png = OUT_PNG / "fig_context_conditions.png"
    fig.savefig(pdf, bbox_inches="tight")
    fig.savefig(png, dpi=300, bbox_inches="tight")
    plt.close(fig)
    print(f"  wrote {pdf}  +  {png}")


def _sub(n: int) -> str:
    """Return the subscript-digit string for a small integer."""
    table = {"0": "₀", "1": "₁", "2": "₂", "3": "₃", "4": "₄",
             "5": "₅", "6": "₆", "7": "₇", "8": "₈", "9": "₉"}
    return "".join(table.get(c, c) for c in str(n))


# ---------------------------------------------------------------------------
# PPTX version (editable)
# ---------------------------------------------------------------------------
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_rect(slide, x, y, w, h, *, fill, line, line_w=1.0, rounded=True):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    if rounded:
        s.adjustments[0] = 0.18
    if fill is not None:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def _add_text(slide, x, y, w, h, text, *, color, size, bold=False,
              italic=False, align="left", vertical="middle", monospace=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if vertical == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif vertical == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = "Consolas" if monospace else "DejaVu Sans"
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return box


def make_pptx() -> None:
    prs = Presentation()
    W, H = 11.0, 6.0
    prs.slide_width = Inches(W); prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    INK_   = _rgb(INK)
    MUTED_ = _rgb(MUTED)
    SUBTLE_ = _rgb(SUBTLE)
    BLUE_  = _rgb(BLUE); BLUE_HI_ = _rgb(BLUE_HI)
    ORANGE_ = _rgb(ORANGE); ORANGE_HI_ = _rgb(ORANGE_HI)
    RED_ = _rgb(RED); DARK_RED_ = _rgb(DARK_RED); RED_HI_ = _rgb(RED_HI)
    LIGHT_ = _rgb("#F5F5F7")

    # Title
    _add_text(slide, 0.20, 0.20, W - 0.40, 0.40,
              "Context conditions",
              color=INK_, size=22, bold=True, align="center")
    _add_text(slide, 0.20, 0.62, W - 0.40, 0.30,
              "Same item, same hidden target — only which prior AD clips reach the answerer changes.",
              color=MUTED_, size=11, italic=True, align="center")

    # Geometry (inches)
    label_x = 0.30
    label_w = 2.40
    timeline_x = label_x + label_w + 0.20
    timeline_w = 6.10
    target_w = 1.20
    target_x = timeline_x + timeline_w + 0.20
    n_clips = 6
    clip_gap = 0.07
    cell_w = (timeline_w - clip_gap * (n_clips - 1)) / n_clips
    row_h = 0.60

    rows = [
        ("(i)  No context", "k = 0", MUTED_, [],
         MUTED_, LIGHT_,
         "Answerer receives no prior AD.",
         "→ 0 clips"),
        ("(ii)  Fixed window", "k = 4", BLUE_,
         [2, 3, 4, 5],
         BLUE_, BLUE_HI_,
         "Answerer receives the most recent k AD clips.",
         "→ 4 clips"),
        ("(iii)  Adaptive", "evidence-selected",
         ORANGE_,
         [1, 4],
         ORANGE_, ORANGE_HI_,
         "Answerer receives only the AD clips selected as evidence at construction time — may skip clips, not always the tail.",
         "→ 2 clips"),
    ]

    row_y_start = 1.30
    row_step = 1.45

    for r, (label, sublabel, sub_color, selected, accent, accent_hi,
            caption, tag) in enumerate(rows):
        y = row_y_start + r * row_step

        # Row labels (left)
        _add_text(slide, label_x, y, label_w, 0.34, label,
                  color=INK_, size=14, bold=True, vertical="top")
        _add_text(slide, label_x, y + 0.36, label_w, 0.26, sublabel,
                  color=sub_color, size=10, italic=True, bold=True,
                  vertical="top")

        # Clip cells
        for i in range(n_clips):
            x = timeline_x + i * (cell_w + clip_gap)
            is_selected = i in selected
            if is_selected:
                _add_rect(slide, x, y, cell_w, row_h,
                          fill=accent_hi, line=accent, line_w=1.6)
                _add_text(slide, x, y, cell_w, row_h,
                          f"AD{_sub(i+1)}",
                          color=accent, size=12, bold=True, align="center",
                          vertical="middle")
            else:
                _add_rect(slide, x, y, cell_w, row_h,
                          fill=LIGHT_, line=SUBTLE_, line_w=0.7)
                _add_text(slide, x, y, cell_w, row_h,
                          f"AD{_sub(i+1)}",
                          color=SUBTLE_, size=11, align="center",
                          vertical="middle")

        # Target — hatched red box with two-line "next AD / (hidden)" label.
        # We approximate the hatch via a slightly darker fill plus a separate
        # diagonal-line shape, since python-pptx hatching support is limited.
        _add_rect(slide, target_x, y, target_w, row_h,
                  fill=RED_HI_, line=RED_, line_w=1.4)
        # Diagonal stripe overlay (3 thin red lines) to suggest hatching.
        from pptx.enum.shapes import MSO_CONNECTOR as _MC
        for off in (0.20, 0.45, 0.70):
            ln = slide.shapes.add_connector(
                _MC.STRAIGHT,
                Inches(target_x + off * target_w),
                Inches(y + 0.03),
                Inches(target_x + (off - 0.20) * target_w),
                Inches(y + row_h - 0.03),
            )
            ln.line.color.rgb = RED_
            ln.line.width = Pt(0.6)
        # Two-line label
        _add_text(slide, target_x, y, target_w, row_h * 0.55,
                  "next AD",
                  color=DARK_RED_, size=12, bold=True, align="center",
                  vertical="bottom")
        _add_text(slide, target_x, y + row_h * 0.45, target_w, row_h * 0.55,
                  "(hidden)",
                  color=DARK_RED_, size=10, italic=True,
                  align="center", vertical="top")

        # Caption (under timeline)
        _add_text(slide, timeline_x, y + row_h + 0.05, timeline_w, 0.30,
                  caption,
                  color=INK_, size=10, italic=True, vertical="top")
        # Tag (right-aligned, accent color, under target)
        _add_text(slide, target_x - 1.20, y + row_h + 0.05,
                  target_w + 1.20, 0.30,
                  tag,
                  color=accent, size=10, italic=True, bold=True,
                  align="right", vertical="top")

    # Footer
    _add_text(slide, 0.30, H - 0.45, W - 0.60, 0.30,
              "The next AD sentence is hidden in every condition; only the "
              "prior subset given to the answerer changes.",
              color=MUTED_, size=10, italic=True, align="center",
              vertical="middle")

    out = OUT_PPTX / "fig_context_conditions.pptx"
    prs.save(out)
    print(f"  wrote {out}")


def main() -> None:
    print("Building context-conditions schematic…")
    make_matplotlib()
    make_pptx()
    print("Done.")


if __name__ == "__main__":
    main()
