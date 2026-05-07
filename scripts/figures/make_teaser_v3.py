"""Build the ForSeeBench teaser figure following the user-provided
``teaser.png`` structure, but with the Charlie / Charlie St. Cloud example
(``good_examples/example_01``).

Layout, top → bottom:
  1. Dark-teal title bar
  2. ``Existing AD-generation pipeline`` strip (light grey)
       video frames → AD Model → AD text bubble
  3. Green-dashed firewall: ``▼ the answerer reads only what is below this line ▼``
  4. Two-column body:
       LEFT  ── Prior AD context with AD₁/AD₂, bridge row highlighted
       RIGHT ── (top) ``Correct next AD (Hidden)`` pink panel with the
                  target frame and the target sentence;
                (bottom) ``Q: What happens next?`` with four options,
                  the correct option marked ✓.

Outputs:
  figures/paper_figures/fig_teaser_v3.{pdf,png}
  figures/paper_figures_pptx/fig_teaser_v3.pptx
"""
from __future__ import annotations

from pathlib import Path

import matplotlib.image as mpimg
import matplotlib.patches as patches
import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path("/jumbo/jinlab/Sally/ForSeeBench")
FE = ROOT / "figure_examples"
OUT_PNG = ROOT / "figures/paper_figures"
OUT_PPTX = ROOT / "figures/paper_figures_pptx"
OUT_PNG.mkdir(parents=True, exist_ok=True)
OUT_PPTX.mkdir(parents=True, exist_ok=True)

# Palette tuned to match the user's teaser.png look
INK = "#1F2933"
MUTED = "#6B7280"
SUBTLE = "#9CA3AF"
TITLE_BAR = "#1B5E5E"   # dark teal title strip
PANEL_GREY = "#F1F1F2"  # light grey for AD-generation pipeline
ADMODEL_TEAL = "#2D7D8C"
ADMODEL_TEAL_HI = "#D9EAEE"
ADTEXT_ORANGE = "#E78B3A"
ADTEXT_HI = "#FCEEDC"
PRIOR_HI = "#FFF3D2"        # yellow tint for bridge row
KW_GREEN = "#1F8B5E"        # highlighted keywords
KW_BLUE  = "#2D6FB6"
GREEN = "#3F8F5C"
DARK_GREEN = "#1E5631"
GREEN_HI = "#DDEFD2"
RED = "#C0413E"
DARK_RED = "#7B1F1C"
RED_HI = "#FBE3E1"
PINK = "#E29093"
BLUE = "#2D6FB6"


# ---------------------------------------------------------------------------
# Charlie example data
# ---------------------------------------------------------------------------
EX_FOLDER = FE / "good_examples/example_01"

# AD lines visible to the answerer (the two adaptive-context clips).
AD_LINES = [
    ("AD₁", "“Now, Charlie strides out of the restaurant, his eyes cast downward.”"),
    ("AD₂", "“As Charlie steps into the street, a fisherman in waterproof bib pants passes him carrying a coil of rope slung over his shoulder.”"),
]
BRIDGE_INDEX = 1  # AD₂ is the bridge

# Hidden target
TARGET_TEXT = "“Charlie stares mesmerized out at the water.”"

# Question and options
QUESTION = "Q:  What happens next?"
OPTIONS = [
    ("A", "Charlie gazes intently at the water.",                       True),
    ("B", "Charlie looks around the street for a lost pet.",            False),
    ("C", "Charlie picks up a piece of trash on the sidewalk.",         False),
    ("D", "Charlie talks to a passerby about the weather.",             False),
]


def _pick_frames(folder: Path, n: int = 3) -> list[Path]:
    """Pick three representative context frames (one per timestamp tag)."""
    out = []
    for tag in ["t0", "tmid", "tend"]:
        for fp in sorted(folder.glob("frames/context_*.jpg")):
            if fp.name.endswith(f"__{tag}.jpg"):
                out.append(fp)
                break
    return out[:n]


def _pick_target_frame(folder: Path) -> Path | None:
    matches = sorted(folder.glob("frames/target__*tmid*.jpg"))
    return matches[0] if matches else None


# ---------------------------------------------------------------------------
# Matplotlib version
# ---------------------------------------------------------------------------
def _arrow(ax, x1, y1, x2, y2, *, color=BLUE, lw=1.6, ls="-"):
    ax.add_patch(FancyArrowPatch(
        (x1, y1), (x2, y2),
        arrowstyle="-|>", mutation_scale=12, lw=lw, color=color,
        linestyle=ls, zorder=4,
    ))


def _round_box(ax, x, y, w, h, *, ec, fc, lw=1.2, r=0.020, zorder=2):
    ax.add_patch(FancyBboxPatch(
        (x, y), w, h,
        boxstyle=f"round,pad=0.0,rounding_size={r}",
        lw=lw, ec=ec, fc=fc, zorder=zorder,
    ))


def _flat_top_box(ax, x, y, w, h, *, ec, fc, lw=1.2, r=0.020):
    """Rounded box that we then visually flatten on the bottom by overlaying
    a small rectangle of the same colour. Used for header strips."""
    _round_box(ax, x, y, w, h, ec=ec, fc=fc, lw=lw, r=r)
    ax.add_patch(patches.Rectangle(
        (x, y), w, h * 0.30, fc=fc, ec=fc, lw=0, zorder=2,
    ))


def make_matplotlib() -> None:
    plt.rcParams.update({
        "font.family": "DejaVu Sans",
        "font.size": 9,
        "pdf.fonttype": 42, "ps.fonttype": 42,
    })

    fig = plt.figure(figsize=(9.5, 5.6))
    fig.patch.set_facecolor("white")
    ax = fig.add_axes([0, 0, 1, 1]); ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    ax.axis("off")

    # ============================================================
    # 1) Title bar (dark teal)
    # ============================================================
    title_h = 0.075
    ax.add_patch(patches.Rectangle(
        (0, 1 - title_h), 1, title_h,
        fc=TITLE_BAR, ec=TITLE_BAR, lw=0, zorder=1,
    ))
    ax.text(0.5, 1 - title_h / 2,
            "ForSeeBench: Prospective Audio-Description QA",
            ha="center", va="center", fontsize=14, fontweight="bold",
            color="white", zorder=3)

    # ============================================================
    # 2) AD-generation pipeline strip (light grey)
    # ============================================================
    pipe_y_top = 1 - title_h - 0.005
    pipe_y_bot = 0.640
    ax.add_patch(patches.Rectangle(
        (0.012, pipe_y_bot), 0.976, pipe_y_top - pipe_y_bot,
        fc=PANEL_GREY, ec=SUBTLE, lw=0.8, zorder=1,
    ))
    ax.text(0.020, pipe_y_top - 0.024,
            "Existing AD-generation pipeline",
            fontsize=8.6, color=INK, fontweight="bold", va="top", zorder=3)

    # 3 video frames
    fr_y = pipe_y_bot + 0.040
    fr_h = pipe_y_top - pipe_y_bot - 0.080
    fr_w = 0.085
    fr_gap = 0.010
    frames = _pick_frames(EX_FOLDER, 3)
    fr_x_start = 0.020
    for i, fp in enumerate(frames):
        x = fr_x_start + i * (fr_w + fr_gap)
        ax.imshow(mpimg.imread(fp), extent=(x, x + fr_w, fr_y, fr_y + fr_h),
                  aspect="auto", zorder=2)
        ax.add_patch(patches.Rectangle(
            (x, fr_y), fr_w, fr_h, fill=False, lw=0.7, ec="#111", zorder=3,
        ))

    # → arrow
    arr1_x = fr_x_start + 3 * (fr_w + fr_gap) - fr_gap
    arr_y = fr_y + fr_h / 2
    _arrow(ax, arr1_x + 0.005, arr_y, arr1_x + 0.052, arr_y,
           color=ADMODEL_TEAL, lw=2.0)

    # AD Model box (teal)
    am_x = arr1_x + 0.060
    am_w = 0.155
    _round_box(ax, am_x, fr_y, am_w, fr_h,
               ec=ADMODEL_TEAL, fc=ADMODEL_TEAL_HI, lw=1.6, r=0.020)
    ax.text(am_x + am_w / 2, fr_y + fr_h * 0.66, "AD Model",
            ha="center", va="center", fontsize=10.5, fontweight="bold",
            color=ADMODEL_TEAL)
    ax.text(am_x + am_w / 2, fr_y + fr_h * 0.30,
            "(MAD-eval, NarrAD,\nAutoAD-Zero)",
            ha="center", va="center", fontsize=7.4, color=ADMODEL_TEAL,
            linespacing=1.30)

    # → arrow
    _arrow(ax, am_x + am_w + 0.005, arr_y, am_x + am_w + 0.052, arr_y,
           color=ADMODEL_TEAL, lw=2.0)

    # AD text bubble (orange)
    at_x = am_x + am_w + 0.060
    at_w = 1 - at_x - 0.020
    _round_box(ax, at_x, fr_y, at_w, fr_h,
               ec=ADTEXT_ORANGE, fc="#FFFFFF", lw=1.4, r=0.020)
    ax.text(at_x + 0.012, fr_y + fr_h - 0.022, "AD text",
            fontsize=10.5, fontweight="bold", color=ADTEXT_ORANGE, va="top")
    # AD lines inside the bubble
    yy = fr_y + fr_h - 0.060
    for tag, line in AD_LINES:
        # truncate display to fit
        ax.text(at_x + 0.012, yy, line,
                fontsize=8.2, color=INK, va="top")
        yy -= 0.054

    # ============================================================
    # 3) Firewall divider (green dashed line)
    # ============================================================
    fw_y = 0.605
    ax.plot([0.012, 0.988], [fw_y, fw_y], color=GREEN,
            lw=1.4, ls=(0, (4, 3)), zorder=3)
    ax.text(0.5, fw_y - 0.024,
            "▼  the answerer reads only what is below this line  ▼",
            ha="center", va="center", fontsize=10, fontweight="bold",
            color=GREEN, style="italic", zorder=4)

    # ============================================================
    # 4) Two-column body
    # ============================================================
    body_y_top = 0.560
    body_y_bot = 0.040
    body_h = body_y_top - body_y_bot

    # ── LEFT: Prior AD context ──
    left_x = 0.012
    left_w = 0.560
    ax.add_patch(FancyBboxPatch(
        (left_x, body_y_bot), left_w, body_h,
        boxstyle="round,pad=0.0,rounding_size=0.014",
        lw=1.4, ec=BLUE, fc="#FFFFFF", zorder=2,
    ))
    # Header strip
    header_h = 0.060
    _flat_top_box(ax, left_x, body_y_top - header_h, left_w, header_h,
                  ec=BLUE, fc=BLUE, lw=0, r=0.014)
    ax.text(left_x + 0.014, body_y_top - header_h / 2,
            "Prior AD context",
            fontsize=10.5, fontweight="bold", color="white", va="center", zorder=4)

    # AD lines below the header
    inner_top = body_y_top - header_h - 0.018
    line_y = inner_top
    line_h = 0.130
    for i, (tag, text) in enumerate(AD_LINES):
        is_bridge = (i == BRIDGE_INDEX)
        if is_bridge:
            # Yellow tint behind the bridge row
            ax.add_patch(FancyBboxPatch(
                (left_x + 0.014, line_y - line_h + 0.010),
                left_w - 0.028, line_h,
                boxstyle="round,pad=0.005,rounding_size=0.014",
                lw=0, fc=PRIOR_HI, zorder=2,
            ))
        # Tag (AD₁, AD₂)
        ax.text(left_x + 0.030, line_y - 0.002, tag,
                fontsize=12, fontweight="bold",
                color=ADMODEL_TEAL,
                family="monospace", va="top", zorder=4)
        # Body text — render bridge with two highlighted phrases inline.
        if is_bridge:
            # We render the whole line as plain INK first, then overlay the
            # highlighted phrases in their own colour with bold weight at
            # approximate x positions. The phrases are short enough that
            # this approximation is visually acceptable.
            body = ("As Charlie steps into the street, a "
                    "fisherman in waterproof bib pants passes him carrying a "
                    "coil of rope slung over his shoulder.")
            import textwrap
            wrapped_lines = textwrap.wrap(body, width=58, break_long_words=False)
            ty = line_y - 0.005
            for wl in wrapped_lines:
                ax.text(left_x + 0.080, ty, wl,
                        fontsize=9.0, color=INK, va="top", zorder=4)
                ty -= 0.035
            # Highlighted keyword chip
            ax.text(left_x + 0.030, line_y - 0.115,
                    "▸ bridge keywords:",
                    fontsize=8.2, color=DARK_RED, fontweight="bold",
                    style="italic", va="top", zorder=4)
            ax.text(left_x + 0.205, line_y - 0.115,
                    "fisherman   ·   waterproof   ·   rope",
                    fontsize=9.4, color=KW_GREEN, fontweight="bold",
                    va="top", zorder=4)
        else:
            ax.text(left_x + 0.080, line_y - 0.005, text,
                    fontsize=9.5, color=INK, va="top", zorder=4)
        line_y -= line_h + 0.010

    # Bottom note inside the left column
    ax.text(left_x + 0.014, body_y_bot + 0.022,
            "→  infer the next narrated visual update from this AD line alone.",
            fontsize=8.2, color=MUTED, style="italic", va="bottom")

    # ── RIGHT (top): Correct next AD (Hidden) ──
    right_x = left_x + left_w + 0.015
    right_w = 1 - right_x - 0.012
    rt_h = 0.218
    rt_y_top = body_y_top
    rt_y_bot = rt_y_top - rt_h
    ax.add_patch(FancyBboxPatch(
        (right_x, rt_y_bot), right_w, rt_h,
        boxstyle="round,pad=0.0,rounding_size=0.014",
        lw=1.4, ec=PINK, fc=RED_HI, zorder=2,
    ))
    # Header strip
    _flat_top_box(ax, right_x, rt_y_top - 0.050, right_w, 0.050,
                  ec=PINK, fc=PINK, lw=0, r=0.014)
    ax.text(right_x + right_w / 2, rt_y_top - 0.025,
            "Correct next AD (Hidden)",
            ha="center", va="center", fontsize=10, fontweight="bold",
            color="white", zorder=4)

    # Target frame on the LEFT of this panel
    target_fp = _pick_target_frame(EX_FOLDER)
    inner_pad = 0.014
    fr_w_t = 0.130
    fr_h_t = rt_h - 0.080
    fr_x_t = right_x + inner_pad
    fr_y_t = rt_y_bot + 0.018
    if target_fp:
        ax.imshow(mpimg.imread(target_fp),
                  extent=(fr_x_t, fr_x_t + fr_w_t, fr_y_t, fr_y_t + fr_h_t),
                  aspect="auto", zorder=3)
        ax.add_patch(patches.Rectangle(
            (fr_x_t, fr_y_t), fr_w_t, fr_h_t,
            fill=False, lw=0.8, ec="#111", zorder=4,
        ))

    # Target text on the RIGHT of the frame
    tx = fr_x_t + fr_w_t + 0.014
    tw = right_x + right_w - tx - inner_pad
    ax.text(tx, rt_y_bot + rt_h - 0.075,
            TARGET_TEXT,
            fontsize=10, color=DARK_RED, va="top", style="italic", zorder=4)

    # ── RIGHT (bottom): Question + options ──
    rb_y_top = rt_y_bot - 0.015
    rb_y_bot = body_y_bot
    rb_h = rb_y_top - rb_y_bot
    ax.add_patch(FancyBboxPatch(
        (right_x, rb_y_bot), right_w, rb_h,
        boxstyle="round,pad=0.0,rounding_size=0.014",
        lw=1.4, ec=GREEN, fc="#FFFFFF", zorder=2,
    ))
    # Question line at top of card
    ax.text(right_x + 0.014, rb_y_top - 0.024,
            QUESTION,
            fontsize=10.5, fontweight="bold", color=INK, va="top", zorder=4)

    opt_y_top = rb_y_top - 0.062
    opt_h = (rb_h - 0.090) / 4
    for j, (letter, text, is_correct) in enumerate(OPTIONS):
        oy_top = opt_y_top - j * opt_h
        oy_bot = oy_top - opt_h
        if is_correct:
            ax.add_patch(FancyBboxPatch(
                (right_x + 0.010, oy_bot + 0.005),
                right_w - 0.020, opt_h - 0.008,
                boxstyle="round,pad=0.005,rounding_size=0.012",
                lw=0, fc=GREEN_HI, zorder=2,
            ))
        # Letter
        ax.text(right_x + 0.022, oy_top - opt_h / 2, f"({letter})",
                fontsize=10, fontweight="bold",
                color=GREEN if is_correct else INK,
                va="center", zorder=4)
        ax.text(right_x + 0.060, oy_top - opt_h / 2, text,
                fontsize=9.0, color=INK, va="center", zorder=4)
        if is_correct:
            ax.text(right_x + right_w - 0.018, oy_top - opt_h / 2, "✓",
                    fontsize=14, color=GREEN, fontweight="bold",
                    ha="right", va="center", zorder=4)

    # Save
    pdf = OUT_PNG / "fig_teaser_v3.pdf"
    png = OUT_PNG / "fig_teaser_v3.png"
    fig.savefig(pdf, bbox_inches="tight")
    fig.savefig(png, dpi=300, bbox_inches="tight")
    plt.close(fig)
    print(f"  wrote {pdf}  +  {png}")


# ---------------------------------------------------------------------------
# PPTX version
# ---------------------------------------------------------------------------
def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_rect(slide, x, y, w, h, *, fill, line, line_w=1.0, rounded=True,
              corner=0.10):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    if rounded:
        s.adjustments[0] = corner
    if fill is not None:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def _add_text(slide, x, y, w, h, text, *, color, size, bold=False,
              italic=False, align="left", vertical="middle", monospace=False,
              line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if vertical == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif vertical == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = "Consolas" if monospace else "DejaVu Sans"
    f.size = Pt(size)
    f.bold = bold; f.italic = italic
    f.color.rgb = color
    return box


def _add_arrow(slide, x1, y1, x2, y2, *, color, width=2.0, dashed=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = line.line._get_or_add_ln()
    tail_end = ln.find(qn("a:tailEnd"))
    if tail_end is None:
        tail_end = etree.SubElement(ln, qn("a:tailEnd"))
    tail_end.set("type", "triangle"); tail_end.set("w", "med"); tail_end.set("len", "med")
    if dashed:
        prst_dash = ln.find(qn("a:prstDash"))
        if prst_dash is None:
            prst_dash = etree.SubElement(ln, qn("a:prstDash"))
        prst_dash.set("val", "dash")
    return line


def make_pptx() -> None:
    prs = Presentation()
    W, H = 13.5, 8.0
    prs.slide_width = Inches(W); prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    INK_ = _rgb(INK); MUTED_ = _rgb(MUTED); SUBTLE_ = _rgb(SUBTLE)
    TITLE_ = _rgb(TITLE_BAR); GREY_ = _rgb(PANEL_GREY)
    AM_TEAL_ = _rgb(ADMODEL_TEAL); AM_TEAL_HI_ = _rgb(ADMODEL_TEAL_HI)
    AT_OR_ = _rgb(ADTEXT_ORANGE); AT_HI_ = _rgb(ADTEXT_HI)
    PRIOR_HI_ = _rgb(PRIOR_HI); KW_ = _rgb(KW_GREEN)
    BLUE_ = _rgb(BLUE)
    GREEN_ = _rgb(GREEN); GREEN_HI_ = _rgb(GREEN_HI); DARK_GREEN_ = _rgb(DARK_GREEN)
    RED_ = _rgb(RED); DARK_RED_ = _rgb(DARK_RED); RED_HI_ = _rgb(RED_HI)
    PINK_ = _rgb(PINK); WHITE_ = _rgb("#FFFFFF")

    # 1) Title bar (full width, dark teal, no rounded corners)
    title_h = 0.55
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(title_h),
    )
    title_bg.fill.solid(); title_bg.fill.fore_color.rgb = TITLE_
    title_bg.line.fill.background()
    title_bg.shadow.inherit = False
    _add_text(slide, 0, 0, W, title_h,
              "ForSeeBench: Prospective Audio-Description QA",
              color=WHITE_, size=22, bold=True, align="center",
              vertical="middle")

    # 2) AD-generation pipeline strip
    pipe_y = title_h + 0.10
    pipe_h = 1.95
    _add_rect(slide, 0.20, pipe_y, W - 0.40, pipe_h,
              fill=GREY_, line=SUBTLE_, line_w=0.8, rounded=False)
    _add_text(slide, 0.35, pipe_y + 0.10, 6.0, 0.30,
              "Existing AD-generation pipeline",
              color=INK_, size=12, bold=True, vertical="top")

    # Three video frames
    fr_y = pipe_y + 0.50
    fr_h = pipe_h - 0.70
    fr_w = 1.45
    fr_gap = 0.10
    fr_x_start = 0.40
    frames = _pick_frames(EX_FOLDER, 3)
    for i, fp in enumerate(frames):
        slide.shapes.add_picture(str(fp),
                                 Inches(fr_x_start + i * (fr_w + fr_gap)),
                                 Inches(fr_y),
                                 Inches(fr_w), Inches(fr_h))

    # arrow → AD Model
    a1_x = fr_x_start + 3 * (fr_w + fr_gap) - fr_gap
    arr_y = fr_y + fr_h / 2
    _add_arrow(slide, a1_x + 0.10, arr_y, a1_x + 0.55, arr_y,
               color=AM_TEAL_, width=2.4)

    # AD Model box
    am_x = a1_x + 0.65
    am_w = 1.85
    _add_rect(slide, am_x, fr_y, am_w, fr_h,
              fill=AM_TEAL_HI_, line=AM_TEAL_, line_w=1.6, corner=0.12)
    _add_text(slide, am_x, fr_y + 0.10, am_w, 0.50, "AD Model",
              color=AM_TEAL_, size=14, bold=True, align="center", vertical="middle")
    _add_text(slide, am_x, fr_y + 0.65, am_w, fr_h - 0.75,
              "(MAD-eval, NarrAD,\nAutoAD-Zero)",
              color=AM_TEAL_, size=10, align="center", vertical="top",
              line_spacing=1.30)

    # arrow → AD text
    _add_arrow(slide, am_x + am_w + 0.10, arr_y, am_x + am_w + 0.55, arr_y,
               color=AM_TEAL_, width=2.4)

    # AD text bubble (orange-bordered white box)
    at_x = am_x + am_w + 0.65
    at_w = W - 0.20 - at_x - 0.10
    _add_rect(slide, at_x, fr_y, at_w, fr_h,
              fill=WHITE_, line=AT_OR_, line_w=1.6, corner=0.10)
    _add_text(slide, at_x + 0.15, fr_y + 0.10, at_w - 0.30, 0.30,
              "AD text",
              color=AT_OR_, size=12, bold=True, vertical="top")
    yy = fr_y + 0.50
    for tag, line in AD_LINES:
        _add_text(slide, at_x + 0.15, yy, at_w - 0.30, 0.45,
                  line, color=INK_, size=10, vertical="top",
                  line_spacing=1.20)
        yy += 0.45

    # 3) Firewall divider
    fw_y = pipe_y + pipe_h + 0.20
    fw_line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.20), Inches(fw_y), Inches(W - 0.20), Inches(fw_y),
    )
    fw_line.line.color.rgb = GREEN_
    fw_line.line.width = Pt(1.6)
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = fw_line.line._get_or_add_ln()
    pd = ln.find(qn("a:prstDash"))
    if pd is None:
        pd = etree.SubElement(ln, qn("a:prstDash"))
    pd.set("val", "dash")
    _add_text(slide, 0, fw_y + 0.05, W, 0.32,
              "▼  the answerer reads only what is below this line  ▼",
              color=GREEN_, size=14, bold=True, italic=True, align="center")

    # 4) Body — two columns
    body_y_top = fw_y + 0.55
    body_y_bot = H - 0.30
    body_h = body_y_bot - body_y_top

    left_x = 0.20
    left_w = 7.80
    right_x = left_x + left_w + 0.20
    right_w = W - right_x - 0.20

    # ── LEFT: Prior AD context ──
    _add_rect(slide, left_x, body_y_top, left_w, body_h,
              fill=WHITE_, line=BLUE_, line_w=1.6, corner=0.05)
    # Header
    _add_rect(slide, left_x, body_y_top, left_w, 0.45,
              fill=BLUE_, line=None, corner=0.20)
    _add_text(slide, left_x + 0.20, body_y_top, left_w - 0.40, 0.45,
              "Prior AD context",
              color=WHITE_, size=14, bold=True, vertical="middle")

    # AD lines
    line_y = body_y_top + 0.65
    line_h = 1.10
    for i, (tag, text) in enumerate(AD_LINES):
        is_bridge = (i == BRIDGE_INDEX)
        if is_bridge:
            _add_rect(slide, left_x + 0.20, line_y - 0.05, left_w - 0.40, line_h + 0.10,
                      fill=PRIOR_HI_, line=None, corner=0.06)
        # Tag (AD₁ / AD₂)
        _add_text(slide, left_x + 0.30, line_y, 0.80, 0.40, tag,
                  color=AM_TEAL_, size=14, bold=True, monospace=True,
                  vertical="top")
        # Body text
        _add_text(slide, left_x + 1.10, line_y, left_w - 1.30, line_h,
                  text, color=INK_, size=12, vertical="top",
                  line_spacing=1.30)
        line_y += line_h + 0.10

    # Bridge keywords callout (below AD₂)
    callout_y = line_y - 0.05
    _add_text(slide, left_x + 0.30, callout_y, 2.20, 0.30,
              "▸ bridge keywords:",
              color=DARK_RED_, size=11, bold=True, italic=True, vertical="top")
    _add_text(slide, left_x + 2.55, callout_y, left_w - 2.85, 0.30,
              "fisherman   ·   waterproof   ·   rope",
              color=KW_, size=12, bold=True, vertical="top")
    _add_text(slide, left_x + 0.30, callout_y + 0.34, left_w - 0.60, 0.30,
              "      →  implies a waterfront setting.",
              color=DARK_RED_, size=11, italic=True, vertical="top")

    # Bottom note
    _add_text(slide, left_x + 0.20, body_y_bot - 0.40, left_w - 0.40, 0.30,
              "→  infer the next narrated visual update from this AD line alone.",
              color=MUTED_, size=11, italic=True, vertical="middle")

    # ── RIGHT TOP: Correct next AD (Hidden) ──
    rt_h = body_h * 0.42
    rt_y_top = body_y_top
    _add_rect(slide, right_x, rt_y_top, right_w, rt_h,
              fill=RED_HI_, line=PINK_, line_w=1.6, corner=0.06)
    _add_rect(slide, right_x, rt_y_top, right_w, 0.40,
              fill=PINK_, line=None, corner=0.20)
    _add_text(slide, right_x, rt_y_top, right_w, 0.40,
              "Correct next AD (Hidden)",
              color=WHITE_, size=12, bold=True, align="center", vertical="middle")

    # Target frame on the LEFT of this panel
    target_fp = _pick_target_frame(EX_FOLDER)
    fr_x_t = right_x + 0.20
    fr_y_t = rt_y_top + 0.55
    fr_w_t = 1.45
    fr_h_t = rt_h - 0.75
    if target_fp:
        slide.shapes.add_picture(str(target_fp),
                                 Inches(fr_x_t), Inches(fr_y_t),
                                 Inches(fr_w_t), Inches(fr_h_t))
    # Target text
    tx = fr_x_t + fr_w_t + 0.20
    tw = right_x + right_w - tx - 0.20
    _add_text(slide, tx, fr_y_t, tw, fr_h_t,
              TARGET_TEXT,
              color=DARK_RED_, size=11, italic=True, vertical="top",
              line_spacing=1.30)

    # ── RIGHT BOTTOM: Question + options ──
    rb_y_top = rt_y_top + rt_h + 0.20
    rb_h = body_y_bot - rb_y_top
    _add_rect(slide, right_x, rb_y_top, right_w, rb_h,
              fill=WHITE_, line=GREEN_, line_w=1.6, corner=0.06)
    _add_text(slide, right_x + 0.20, rb_y_top + 0.10, right_w - 0.40, 0.35,
              QUESTION,
              color=INK_, size=12, bold=True, vertical="top")

    opt_top = rb_y_top + 0.55
    opt_h = (rb_h - 0.65) / 4
    for j, (letter, text, is_correct) in enumerate(OPTIONS):
        oy = opt_top + j * opt_h
        if is_correct:
            _add_rect(slide, right_x + 0.15, oy, right_w - 0.30, opt_h - 0.05,
                      fill=GREEN_HI_, line=None, corner=0.10)
        _add_text(slide, right_x + 0.25, oy, 0.50, opt_h - 0.05, f"({letter})",
                  color=GREEN_ if is_correct else INK_, size=11, bold=True,
                  vertical="middle")
        _add_text(slide, right_x + 0.75, oy, right_w - 1.20, opt_h - 0.05,
                  text, color=INK_, size=10, vertical="middle")
        if is_correct:
            _add_text(slide, right_x + right_w - 0.45, oy,
                      0.30, opt_h - 0.05,
                      "✓", color=GREEN_, size=14, bold=True,
                      align="right", vertical="middle")

    out = OUT_PPTX / "fig_teaser_v3.pptx"
    prs.save(out)
    print(f"  wrote {out}")


def main() -> None:
    print("Building Charlie teaser (v3)…")
    make_matplotlib()
    make_pptx()
    print("Done.")


if __name__ == "__main__":
    main()
