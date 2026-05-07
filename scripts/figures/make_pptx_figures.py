"""Build editable PowerPoint versions of the ForSeeBench paper figures.

Outputs (.pptx) → `figures/paper_figures_pptx/`. One slide per figure, sized to
match the matplotlib versions. Every panel, header, text block, frame, and
arrow is its own pptx shape so the user can click and edit anything inside
PowerPoint or Keynote.

Slides written:
  - fig_teaser.pptx
  - fig_filter.pptx
  - fig_ad_source.pptx
  - fig_good_gallery.pptx
  - fig_ad_source_grid.pptx
  - fig_target_types.pptx
  - fig_evidence_chain.pptx
  - all_figures.pptx              (one slide per figure, single deck)
"""
from __future__ import annotations

import json
import re
import textwrap
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path("/jumbo/jinlab/Sally/ForSeeBench")
FE = ROOT / "figure_examples"
OUT = ROOT / "figures/paper_figures_pptx"
OUT.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------------------
# Palette — same as matplotlib figures, encoded as RGB tuples
# ---------------------------------------------------------------------------
def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


INK       = _rgb("#1F2933")
MUTED     = _rgb("#6B7280")
SUBTLE    = _rgb("#9CA3AF")
WHITE     = _rgb("#FFFFFF")
BLACK     = _rgb("#111111")
BLUE      = _rgb("#3D6FB6")
DARK_BLUE = _rgb("#1E3A8A")
GREEN     = _rgb("#3F8F5C")
DARK_GREEN= _rgb("#1E5631")
RED       = _rgb("#C0413E")
DARK_RED  = _rgb("#7B1F1C")
ORANGE    = _rgb("#D97A2D")
GOLD      = _rgb("#E2A83C")
PURPLE    = _rgb("#7E5BB0")
TEAL      = _rgb("#1B7F8E")
PRIOR_HI  = _rgb("#FFF3D2")
GOLD_HI   = _rgb("#DDEFD2")
RED_HI    = _rgb("#FBE3E1")
BLUE_HI   = _rgb("#E8F0FB")
PURPLE_HI = _rgb("#EFE7F4")
PANEL_FACE = _rgb("#FFFFFF")
PANEL_TINT = _rgb("#F8FAFC")
TEAL_HI   = _rgb("#DCEEEC")
ORANGE_HI = _rgb("#FCEEDB")


# ---------------------------------------------------------------------------
# Helpers — measurement and shape construction
# ---------------------------------------------------------------------------
def in_(v: float) -> Emu:
    return Inches(v)


def set_solid_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color: RGBColor | None = None, width: float = 1.0) -> None:
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
        shape.line.width = Pt(width)


def add_rect(slide, x: float, y: float, w: float, h: float, *,
             fill=PANEL_FACE, line=BLUE, line_w=1.4,
             rounded: bool = True):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        in_(x), in_(y), in_(w), in_(h),
    )
    if rounded:
        # Adjust corner radius to a small fraction of shape size.
        s.adjustments[0] = 0.06
    if fill is not None:
        set_solid_fill(s, fill)
    else:
        s.fill.background()
    set_line(s, line, line_w) if line is not None else set_line(s, None)
    s.shadow.inherit = False
    return s


def add_text(slide, x: float, y: float, w: float, h: float, text: str, *,
             color=INK, size: float = 9.0, bold: bool = False,
             italic: bool = False, align: str = "left",
             vertical: str = "top", font: str = "DejaVu Sans",
             monospace: bool = False, line_spacing: float | None = None):
    box = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if vertical == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif vertical == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif vertical == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    tf.text = ""  # start fresh
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = "Consolas" if monospace else font
    f.size = Pt(size)
    f.bold = bool(bold)
    f.italic = bool(italic)
    f.color.rgb = color
    return box


def add_multitext(slide, x, y, w, h, runs: list[dict], *,
                  align="left", vertical="top", line_spacing=None):
    """A single text frame with multiple runs (different styles / colors).
    Each run dict supports: text, color, size, bold, italic, monospace, newline.
    """
    box = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if vertical == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif vertical == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    if line_spacing is not None:
        p.line_spacing = line_spacing

    first = True
    for spec in runs:
        if spec.get("newline"):
            p = tf.add_paragraph()
            p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                           "right": PP_ALIGN.RIGHT}[align]
            if line_spacing is not None:
                p.line_spacing = line_spacing
            first = True
            continue
        text = spec.get("text", "")
        if not text:
            continue
        if first:
            run = p.add_run()
            first = False
        else:
            run = p.add_run()
        run.text = text
        f = run.font
        f.name = "Consolas" if spec.get("monospace") else "DejaVu Sans"
        f.size = Pt(spec.get("size", 9.0))
        f.bold = bool(spec.get("bold"))
        f.italic = bool(spec.get("italic"))
        f.color.rgb = spec.get("color", INK)
    return box


def add_panel(slide, x, y, w, h, *, ec=BLUE, fc=PANEL_FACE,
              header: str | None = None, header_color: RGBColor | None = None,
              header_size: float = 11.0, header_height: float = 0.36):
    """Add a rounded rectangle panel with optional header bar.

    Returns the panel shape (the body) and the header shape (or None).
    """
    body = add_rect(slide, x, y, w, h, fill=fc, line=ec, line_w=1.6)
    header_shape = None
    if header:
        # Header bar at the top — full-width rounded; sits on top of body.
        bar = add_rect(slide, x, y, w, header_height,
                       fill=header_color or ec, line=None, line_w=0)
        # We want the bottom of the header to be flat. Cover with a small
        # rectangle so corners are square at the bottom.
        cover = add_rect(slide, x, y + header_height - 0.08, w, 0.10,
                         fill=header_color or ec, line=None, line_w=0,
                         rounded=False)
        # Header text (white, bold)
        add_text(slide, x + 0.10, y + 0.025, w - 0.20, header_height - 0.05,
                 header, color=WHITE, size=header_size, bold=True,
                 vertical="middle")
        header_shape = bar
    return body, header_shape


def add_arrow(slide, x1, y1, x2, y2, *, color=BLUE, width: float = 2.0,
              dashed: bool = False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, in_(x1), in_(y1), in_(x2), in_(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    # Add arrowhead by tweaking xml
    from pptx.oxml.ns import qn
    ln = line.line._get_or_add_ln()
    tail_end = ln.find(qn("a:tailEnd"))
    if tail_end is None:
        from lxml import etree
        tail_end = etree.SubElement(ln, qn("a:tailEnd"))
    tail_end.set("type", "triangle")
    tail_end.set("w", "med")
    tail_end.set("len", "med")
    if dashed:
        prst_dash = ln.find(qn("a:prstDash"))
        if prst_dash is None:
            from lxml import etree
            prst_dash = etree.SubElement(ln, qn("a:prstDash"))
        prst_dash.set("val", "dash")
    return line


def add_image(slide, path: Path, x: float, y: float, w: float, h: float):
    """Embed an image at the given EMU rect."""
    return slide.shapes.add_picture(str(path), in_(x), in_(y), in_(w), in_(h))


# ---------------------------------------------------------------------------
# Data loaders (mirror make_paper_figures.py / make_extra_figures.py)
# ---------------------------------------------------------------------------
def load_meta(folder: Path) -> dict:
    return json.loads((folder / "metadata.json").read_text())


def first_frame(folder: Path, pattern: str = "context_01*tmid*") -> Path | None:
    matches = sorted(folder.glob(f"frames/{pattern}.jpg"))
    return matches[0] if matches else None


def all_context_frames(folder: Path, pattern: str = "context_*.jpg") -> list[Path]:
    return sorted(folder.glob(f"frames/{pattern}"))


def short_movie(movie: str, limit: int = 30) -> str:
    if len(movie) <= limit:
        return movie
    body = movie
    if re.match(r"^\d{4}_", body):
        body = body[5:]
    pretty = body.replace("_", " ")
    if len(pretty) <= limit:
        return pretty
    return pretty[:limit - 1].rstrip() + "…"


# ---------------------------------------------------------------------------
# Slide-level builders, each producing one slide on a presentation
# ---------------------------------------------------------------------------
def _new_slide(prs: Presentation, w_in: float, h_in: float):
    """Configure the deck's slide size and add a blank slide."""
    prs.slide_width = in_(w_in)
    prs.slide_height = in_(h_in)
    blank = prs.slide_layouts[6]  # blank layout
    return prs.slides.add_slide(blank)


def build_teaser(prs: Presentation):
    """Tracy / The Roommate teaser. Two-region layout:

      TOP: AD-generation pipeline (video frames → AD model → AD text) on a
           grey background, with a green annotation "only the AD text
           crosses to the answerer."
      DIVIDER: a horizontal "firewall" line.
      BOTTOM: answerer's view — prior AD context, hidden target ???, and
              the four-option question.
    """
    W, H = 13.5, 9.0
    slide = _new_slide(prs, W, H)

    # ---- Title bar ----
    add_text(slide, 0.20, 0.18, W - 0.40, 0.40,
             "ForSeeBench: Prospective Audio-Description QA",
             color=INK, size=22, bold=True, align="center")
    add_text(slide, 0.20, 0.62, W - 0.40, 0.30,
             "The answerer is text-only — it reads prior AD lines, never the video.",
             color=MUTED, size=12, italic=True, align="center")

    # ============================================================
    # TOP REGION — AD-generation pipeline (grey background)
    # ============================================================
    top_y0, top_y1 = 1.10, 3.05
    add_rect(slide, 0.30, top_y0, W - 0.60, top_y1 - top_y0,
             fill=_rgb("#F5F5F7"), line=SUBTLE, line_w=1.0)
    add_text(slide, 0.45, top_y0 + 0.06, 7.0, 0.28,
             "AD-generation pipeline  ·  the answerer never sees this",
             color=MUTED, size=10, italic=True, bold=True, vertical="top")

    # 3 video frames (Tracy bedroom, full intensity)
    ex = FE / "baseline_comparison_examples/example_01"
    ctx_frames = all_context_frames(ex, "context_*.jpg")
    chosen = []
    for tag in ["t0", "tmid", "tend"]:
        for fp in ctx_frames:
            if fp.name.endswith(f"__{tag}.jpg"):
                chosen.append(fp)
                break
    chosen = chosen[:3]

    fr_y = top_y0 + 0.50
    fr_h = 1.30
    fr_w = 1.40
    fr_gap = 0.12
    fr_x_start = 0.50
    for i, fp in enumerate(chosen):
        add_image(slide, fp, fr_x_start + i * (fr_w + fr_gap), fr_y, fr_w, fr_h)
    add_text(slide, fr_x_start, top_y1 - 0.30, 3 * (fr_w + fr_gap), 0.22,
             "video frames",
             color=MUTED, size=9, italic=True, align="center")

    # Arrow → AD model
    arrow1_x = fr_x_start + 3 * (fr_w + fr_gap)
    arrow1_y = fr_y + fr_h / 2
    add_arrow(slide, arrow1_x, arrow1_y, arrow1_x + 0.40, arrow1_y,
              color=BLUE, width=2.4)

    # AD model box
    model_x = arrow1_x + 0.45
    model_w = 1.55
    add_rect(slide, model_x, fr_y, model_w, fr_h,
             fill=BLUE_HI, line=BLUE, line_w=1.5)
    add_text(slide, model_x, fr_y, model_w, fr_h,
             "AD model\n(MAD-eval, NarrAD,\nAutoAD-Zero)",
             color=DARK_BLUE, size=11, bold=True, align="center",
             vertical="middle")

    # Arrow → AD text
    add_arrow(slide, model_x + model_w, arrow1_y,
              model_x + model_w + 0.40, arrow1_y,
              color=BLUE, width=2.4)

    # AD text bubble
    text_x = model_x + model_w + 0.45
    text_w = W - 0.30 - text_x - 0.15
    add_rect(slide, text_x, fr_y, text_w, fr_h,
             fill=PANEL_FACE, line=ORANGE, line_w=1.2)
    add_text(slide, text_x + 0.12, fr_y + 0.10, text_w - 0.24, 0.32,
             "AD text",
             color=ORANGE, size=11, bold=True, vertical="top")
    add_text(slide, text_x + 0.12, fr_y + 0.45, text_w - 0.24, fr_h - 0.55,
             "“In her room, Tracy sits up\nin bed and rubs her brow.”",
             color=INK, size=10.5, vertical="top")

    # ============================================================
    # FIREWALL DIVIDER
    # ============================================================
    fw_y = top_y1 + 0.25
    # Use a dotted-style green line via a connector
    fw_line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(0.30), Inches(fw_y),
        Inches(W - 0.30), Inches(fw_y),
    )
    fw_line.line.color.rgb = GREEN
    fw_line.line.width = Pt(1.5)
    # dashed
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = fw_line.line._get_or_add_ln()
    prst_dash = ln.find(qn("a:prstDash"))
    if prst_dash is None:
        prst_dash = etree.SubElement(ln, qn("a:prstDash"))
    prst_dash.set("val", "dash")

    add_text(slide, 0.30, fw_y + 0.05, W - 0.60, 0.30,
             "▼  the answerer reads only what is below this line  ▼",
             color=GREEN, size=12, bold=True, italic=True, align="center")

    # ============================================================
    # BOTTOM REGION — answerer's view
    # ============================================================
    body_top = fw_y + 0.55
    body_bot = H - 0.65

    # Column boundaries
    col_b = (0.30, 8.50)         # prior AD
    col_c = (8.70, W - 0.30)     # right stack

    # ---- Prior AD context (left) ----
    b_x, b_w = col_b[0], col_b[1] - col_b[0]
    add_rect(slide, b_x, body_top, b_w, body_bot - body_top,
             fill=PANEL_FACE, line=BLUE, line_w=1.4)
    # Header strip
    add_rect(slide, b_x, body_top, b_w, 0.42,
             fill=BLUE, line=None, rounded=True)
    add_text(slide, b_x + 0.15, body_top + 0.04, b_w - 0.30, 0.35,
             "Prior AD context  ·  visible to the answerer",
             color=WHITE, size=12, bold=True, vertical="middle")

    # AD₁ — single bridge line, with soft tint
    ad1_y = body_top + 0.65
    add_rect(slide, b_x + 0.15, ad1_y, b_w - 0.30, 1.10,
             fill=PRIOR_HI, line=None)
    add_text(slide, b_x + 0.25, ad1_y + 0.10, 0.50, 0.32, "AD₁",
             color=ORANGE, size=12, bold=True, monospace=True)
    add_text(slide, b_x + 0.78, ad1_y + 0.10, 0.85, 0.32, "−00:08",
             color=ORANGE, size=11.5, bold=True, monospace=True)
    add_text(slide, b_x + 1.70, ad1_y + 0.10, b_w - 1.85, 0.85,
             "In her room, Tracy sits up in bed and rubs her brow.",
             color=INK, size=12.5, line_spacing=1.30)

    # Bridge cue callout
    bridge_y = ad1_y + 1.30
    add_text(slide, b_x + 0.25, bridge_y, 1.50, 0.30,
             "▸ bridge cue",
             color=DARK_RED, size=10.5, bold=True, italic=True)
    add_text(slide, b_x + 1.85, bridge_y, b_w - 2.00, 0.30,
             "“sits up in bed”   +   “rubs her brow”",
             color=ORANGE, size=11.5, bold=True)
    add_text(slide, b_x + 0.25, bridge_y + 0.36, b_w - 0.40, 0.30,
             "      →  she is waking up;  about to leave the bed.",
             color=DARK_RED, size=10.5, italic=True)

    # Bottom note
    add_text(slide, b_x, body_bot - 0.45, b_w, 0.30,
             "→ infer the next narrated visual update from this AD line alone.",
             color=MUTED, size=10.5, italic=True, align="center")

    # ---- Right stack: hidden ??? + QA ----
    c_x, c_w = col_c[0], col_c[1] - col_c[0]
    body_h = body_bot - body_top
    c_top_h = body_h * 0.32
    c_bot_h = body_h * 0.65
    c_top_y0 = body_top
    c_top_y1 = c_top_y0 + c_top_h
    c_bot_y0 = c_top_y1 + body_h * 0.03
    c_bot_y1 = body_bot

    # Hidden future AD card
    add_rect(slide, c_x, c_top_y0, c_w, c_top_h,
             fill=RED_HI, line=RED, line_w=1.4)
    add_text(slide, c_x, c_top_y0 + 0.08, c_w, 0.32,
             "Future AD  ·  hidden",
             color=DARK_RED, size=11, bold=True, align="center")
    add_text(slide, c_x, c_top_y0 + c_top_h / 2 - 0.50, c_w, 1.0,
             "?  ?  ?",
             color=RED, size=44, bold=True, align="center", vertical="middle")
    add_text(slide, c_x + 0.10, c_top_y1 - 0.45, c_w - 0.20, 0.35,
             "target:  “Tracy groggily gets up and crosses the dim room.”",
             color=MUTED, size=8, italic=True, align="center")

    # QA card
    add_rect(slide, c_x, c_bot_y0, c_w, c_bot_y1 - c_bot_y0,
             fill=PANEL_FACE, line=GREEN, line_w=1.4)
    add_text(slide, c_x + 0.15, c_bot_y0 + 0.08, c_w - 0.30, 0.32,
             "Q:  What happens next?",
             color=INK, size=12, bold=True)

    options = [
        ("A", "Tracy gets out of bed and walks across the room.", True),
        ("B", "Tracy sits back down on the bed and continues to rub her brow.", False),
        ("C", "Tracy picks up a book and starts reading.", False),
        ("D", "Tracy talks to someone off-screen.", False),
    ]
    opt_h = 0.55
    opt_gap = 0.05
    opt_top = c_bot_y0 + 0.55
    for j, (letter, opt, is_gold) in enumerate(options):
        oy = opt_top + j * (opt_h + opt_gap)
        if is_gold:
            add_rect(slide, c_x + 0.10, oy, c_w - 0.20, opt_h,
                     fill=GOLD_HI, line=None)
        add_text(slide, c_x + 0.18, oy, 0.45, opt_h, f"({letter})",
                 color=GREEN if is_gold else INK, size=11, bold=True,
                 vertical="middle")
        add_text(slide, c_x + 0.62, oy, c_w - 0.95, opt_h, opt,
                 color=INK, size=10, vertical="middle")
        if is_gold:
            add_text(slide, c_x + c_w - 0.40, oy, 0.30, opt_h, "✓",
                     color=GREEN, size=14, bold=True,
                     align="right", vertical="middle")

    # ---- Footer ----
    add_text(slide, 0.30, H - 0.45, W - 0.60, 0.30,
             "All four options are independently plausible bedroom actions.  "
             "Only the bridge cue (sits up + rubs brow → waking up) makes "
             "“gets out of bed and walks across the room” the natural pick.",
             color=INK, size=10, italic=True, align="center")


def _qa_panel_pptx(slide, x, y, w, h, *, header, header_color, badge,
                   ctx_lines, target_text, options, correct_idx, distractors,
                   frame_path, takeaway, takeaway_color):
    add_panel(slide, x, y, w, h, ec=header_color, fc=PANEL_FACE,
              header=header, header_color=header_color, header_size=10)
    # Badge in header
    add_text(slide, x + w - 1.20, y + 0.06, 1.10, 0.30,
             badge, color=WHITE, size=10, bold=True,
             align="right", vertical="middle")

    # Frame thumbnail
    fx, fy, fw, fh = x + 0.20, y + 0.50, 1.40, 1.05
    if frame_path and frame_path.exists():
        add_image(slide, frame_path, fx, fy, fw, fh)
    else:
        add_rect(slide, fx, fy, fw, fh, fill=_rgb("#E5E7EB"),
                 line=_rgb("#9CA3AF"), line_w=0.6, rounded=False)
        add_text(slide, fx, fy, fw, fh,
                 "(raw clip not\nredistributed)",
                 color=MUTED, size=9, align="center", vertical="middle")
    add_text(slide, fx, fy + fh + 0.02, fw, 0.20,
             "context frame", color=MUTED, size=8, align="center")

    # Prior AD on the right of the frame
    tx = fx + fw + 0.20
    tw = (x + w) - tx - 0.20
    add_text(slide, tx, y + 0.50, tw, 0.30,
             "Prior AD context", color=MUTED, size=10, bold=True)
    yy = y + 0.78
    for line in ctx_lines:
        add_text(slide, tx, yy, 0.20, 0.30, "•",
                 color=INK, size=10, bold=True)
        rows = max(1, (len(line) + 49) // 50)
        add_text(slide, tx + 0.20, yy, tw - 0.20,
                 0.32 * rows, line, color=INK, size=10)
        yy += 0.32 * rows + 0.10

    # Hidden target
    add_text(slide, x + 0.20, y + h - 2.65, 1.80, 0.30,
             "Hidden target AD:", color=DARK_RED, size=10, bold=True)
    add_text(slide, x + 1.95, y + h - 2.65, w - 2.20, 0.55,
             target_text, color=INK, size=10, italic=True)

    # Options
    yy = y + h - 2.10
    for j, opt in enumerate(options):
        letter = chr(ord("A") + j)
        is_gold = (j == correct_idx)
        if is_gold:
            add_rect(slide, x + 0.15, yy - 0.05, w - 0.30, 0.45,
                     fill=GOLD_HI, line=None)
        add_text(slide, x + 0.30, yy, 0.45, 0.40, f"({letter})",
                 color=GREEN if is_gold else INK, size=11, bold=True)
        add_text(slide, x + 0.80, yy, w - 2.50, 0.45, opt,
                 color=INK, size=10)
        if is_gold:
            add_text(slide, x + w - 1.45, yy, 0.40, 0.30, "✓",
                     color=GREEN, size=12, bold=True, align="right")
        add_text(slide, x + w - 1.95, yy + 0.22, 1.75, 0.20,
                 distractors[j], color=MUTED, size=7,
                 monospace=True, align="right")
        yy += 0.42

    # Takeaway tag
    add_rect(slide, x + 0.15, y + h - 0.55, w - 0.30, 0.45,
             fill=takeaway_color if isinstance(takeaway_color, RGBColor) else takeaway_color,
             line=None)
    # The fill above is too saturated; use the lighter HI version. Redo:
    # (We pass the HI color directly via `takeaway_color`.)

    add_text(slide, x + 0.20, y + h - 0.50, w - 0.40, 0.40,
             takeaway, color=INK, size=9, italic=True, bold=True,
             align="center", vertical="middle")


def build_filter(prs: Presentation):
    W, H = 13.0, 8.5
    slide = _new_slide(prs, W, H)
    add_text(slide, 0.20, 0.10, W - 0.40, 0.40,
             "Why nontriviality filtering matters",
             color=INK, size=22, bold=True, align="center")
    add_text(slide, 0.20, 0.52, W - 0.40, 0.35,
             "Construction-time filters drop items answerable without prior-AD evidence. "
             "Retained items keep the typed-distractor structure that requires forward-relevant context.",
             color=MUTED, size=11, italic=True, align="center")

    # Column headers
    add_text(slide, 0.30, 0.97, (W - 0.60) / 2 - 0.05, 0.30,
             "REJECTED  (validation: distractor_quality = 0.0)",
             color=RED, size=14, bold=True, align="center")
    add_text(slide, W / 2 + 0.05, 0.97, (W - 0.60) / 2 - 0.05, 0.30,
             "RETAINED  (787-item benchmark)",
             color=GREEN, size=14, bold=True, align="center")

    panel_w = (W - 0.6) / 2
    panel_h = 3.4
    top_y = 1.40
    bot_y = top_y + panel_h + 0.30

    # Top-left bad_04
    bad4 = FE / "bad_examples/example_04"
    bad4_meta = load_meta(bad4)
    bad4_frame = first_frame(bad4, "context_02*tmid*")
    _qa_panel_pptx(
        slide, 0.30, top_y, panel_w, panel_h,
        header="bad_04 · 1027_Les_Miserables · participant_update",
        header_color=RED, badge="REJECTED",
        ctx_lines=[
            "VALJEAN gets up and hauls MARIUS up onto his shoulder.",
            "With the unconscious MARIUS over his shoulder, VALJEAN wades up through the knee-deep sewerage.",
        ],
        target_text="VALJEAN holds MARIUS up next to him.",
        options=[
            "VALJEAN supports MARIUS while navigating the sewer.",
            "MARIUS regains consciousness and starts walking.",
            "VALJEAN drops MARIUS and runs away.",
            "VALJEAN throws MARIUS into the sewer.",
        ],
        correct_idx=0, distractors=bad4_meta["distractor_metadata"],
        frame_path=bad4_frame,
        takeaway="(A) is a near-paraphrase of both the target and the prior AD → "
                 "answerable from wording alone, no forward evidence required.",
        takeaway_color=RED_HI,
    )

    # Bottom-left bad_03
    bad3 = FE / "bad_examples/example_03"
    bad3_meta = load_meta(bad3)
    _qa_panel_pptx(
        slide, 0.30, bot_y, panel_w, panel_h,
        header="bad_03 · 1026_Legion · spatial_consequence",
        header_color=RED, badge="REJECTED",
        ctx_lines=[
            "MICHAEL takes a nozzle from a fuel pump, sprays gas over the back of the van, and lights it.",
        ],
        target_text="The possessed flee as the van is engulfed in flames.",
        options=[
            "The van bursts into flames.",
            "The possessed chase Michael away.",
            "The van explodes.",
            "Michael runs away from the van.",
        ],
        correct_idx=0, distractors=bad3_meta["distractor_metadata"],
        frame_path=None,
        takeaway="Target is the immediate physical consequence of the prior AD (lights it → bursts into flames). "
                 "Options (A) and (C) are near-synonyms.",
        takeaway_color=RED_HI,
    )

    # Top-right good_03
    g3 = FE / "good_examples/example_03"
    g3_meta = load_meta(g3)
    g3_frame = first_frame(g3, "context_01*tmid*")
    _qa_panel_pptx(
        slide, W / 2 + 0.05, top_y, panel_w, panel_h,
        header="good_03 · 1051_Harry_Potter… · state_change",
        header_color=GREEN, badge="KEPT",
        ctx_lines=["VOLDEMORT forces HARRY to bow."],
        target_text="HARRY writhes in agony.",
        options=[
            "HARRY experiences intense pain.",
            "HARRY starts to laugh uncontrollably.",
            "HARRY suddenly gains superhuman strength.",
            "HARRY begins to float upwards.",
        ],
        correct_idx=0, distractors=g3_meta["distractor_metadata"],
        frame_path=g3_frame,
        takeaway="Prior AD provides a causal precondition (forced to bow) for the future state. "
                 "Distractors are concrete and unambiguously wrong.",
        takeaway_color=GOLD_HI,
    )

    # Bottom-right good_02
    g2 = FE / "good_examples/example_02"
    g2_meta = load_meta(g2)
    g2_frame = first_frame(g2, "context_02*tmid*")
    _qa_panel_pptx(
        slide, W / 2 + 0.05, bot_y, panel_w, panel_h,
        header="good_02 · 3031_HANSEL_GRETEL… · participant_update",
        header_color=GREEN, badge="KEPT",
        ctx_lines=["Mina moves towards Hansel.", "Mina kisses him."],
        target_text="They embrace passionately.",
        options=[
            "Mina and Hansel kiss passionately.",
            "Hansel and Gretel kiss passionately.",
            "The witch appears and chases them away.",
            "The forest becomes covered in snow.",
        ],
        correct_idx=0, distractors=g2_meta["distractor_metadata"],
        frame_path=g2_frame,
        takeaway="Named participants in prior AD make (B) a clean entity_swapped distractor. "
                 "(C) and (D) probe movie-prior shortcuts.",
        takeaway_color=GOLD_HI,
    )


def build_ad_source(prs: Presentation):
    W, H = 13.5, 8.0
    slide = _new_slide(prs, W, H)
    add_text(slide, 0.20, 0.10, W - 0.40, 0.40,
             "AD-source substitution: same item, three AD sources, three predictions",
             color=INK, size=22, bold=True, align="center")
    add_text(slide, 0.20, 0.55, W - 0.40, 0.35,
             "ForSeeBench fixes the question, options, and selected context positions; "
             "only the AD source filling those positions changes "
             "(item: baseline_comparison_examples/example_01).",
             color=MUTED, size=11, italic=True, align="center")

    # Top panel: shared item
    add_panel(slide, 0.30, 1.05, W - 0.60, 2.30,
              ec=BLUE, fc=PANEL_FACE,
              header="Shared item   ·   selected adaptive context: 1 clip",
              header_color=BLUE, header_size=11)

    bc1 = FE / "baseline_comparison_examples/example_01"
    ctx_frames = all_context_frames(bc1, "context_01*.jpg")
    for i, fp in enumerate(ctx_frames[:3]):
        add_image(slide, fp, 0.50 + i * 1.30, 1.55, 1.15, 1.40)
    for i, label in enumerate(["t=0", "t=mid", "t=end"]):
        add_text(slide, 0.50 + i * 1.30, 2.97, 1.15, 0.20, label,
                 color=MUTED, size=9, align="center")
    add_text(slide, 0.50, 3.18, 3.6, 0.18,
             "context frames (visualization only — answerer is text-only)",
             color=MUTED, size=8, italic=True, align="center")

    # Question + options (right)
    add_text(slide, 4.60, 1.50, 5.5, 0.30,
             "Q:  What happens next?",
             color=INK, size=12, bold=True)
    options = [
        "Tracy gets out of bed and walks across the room.",
        "Tracy sits back down on the bed and continues to rub her brow.",
        "Tracy picks up a book and starts reading.",
        "Tracy talks to someone off-screen.",
    ]
    yy = 1.85
    for j, opt in enumerate(options):
        letter = chr(ord("A") + j)
        is_gold = (j == 0)
        if is_gold:
            add_rect(slide, 4.55, yy - 0.04, 5.40, 0.36,
                     fill=GOLD_HI, line=None)
        add_text(slide, 4.65, yy, 0.40, 0.30, f"({letter})",
                 color=GREEN if is_gold else INK, size=10, bold=True)
        add_text(slide, 5.10, yy, 4.50, 0.30, opt,
                 color=INK, size=10)
        yy += 0.34

    # Hidden target on far right
    add_panel(slide, 10.20, 1.50, W - 10.50, 1.85,
              ec=RED, fc=RED_HI,
              header="Hidden target AD", header_color=RED,
              header_size=10, header_height=0.30)
    add_text(slide, 10.20, 2.00, W - 10.50, 0.40,
             "?  ?", color=RED, size=28, bold=True, align="center")
    add_text(slide, 10.30, 2.45, W - 10.70, 0.55,
             "“Tracy groggily gets up and crosses the dim room.”",
             color=DARK_RED, size=9, italic=True, align="center")
    add_text(slide, 10.20, 3.05, W - 10.50, 0.20,
             "withheld during evaluation",
             color=MUTED, size=8, italic=True, align="center")

    # Three rows for AD sources
    rows = [
        ("Human MAD-eval AD",
         "In her room, Tracy sits up in bed and rubs her brow.",
         True,
         "Tracy gets out of bed and walks across the room.",
         GREEN, GOLD_HI,
         "preserves the entity Tracy and the location bed / room",
         "kept_keywords:  Tracy · bed · room"),
        ("NarrAD",
         "Rebecca watches Sarah sleep, then wipes her face thoughtfully.",
         False,
         "Tracy sits back down on the bed and continues to rub her brow.",
         RED, RED_HI,
         "drops Tracy → wrong protagonist (Rebecca / Sarah)",
         "dropped_keywords:  Tracy"),
        ("AutoAD-Zero",
         "Tracy covers her face with her hands in front of the mirror.",
         False,
         "Tracy sits back down on the bed and continues to rub her brow.",
         RED, RED_HI,
         "drops bed → wrong location (mirror)",
         "dropped_keywords:  bed · room"),
    ]
    base_y = 3.55
    row_h = 1.30
    row_gap = 0.12
    for i, (label, ad_text, is_correct, pred, color, hi_color,
            drop_note, kw_note) in enumerate(rows):
        ry = base_y + i * (row_h + row_gap)
        add_panel(slide, 0.30, ry, W - 0.60, row_h,
                  ec=color, fc=hi_color, header=None)
        # Source label badge
        add_rect(slide, 0.45, ry + 0.20, 2.0, 0.85,
                 fill=color, line=None)
        add_text(slide, 0.45, ry + 0.20, 2.0, 0.85, label,
                 color=WHITE, size=12, bold=True, align="center", vertical="middle")
        # AD text
        add_text(slide, 2.65, ry + 0.10, 6.0, 0.25,
                 "AD filling the selected context position:",
                 color=MUTED, size=9, bold=True)
        add_text(slide, 2.65, ry + 0.32, 6.0, 0.45, ad_text,
                 color=INK, size=10.5)
        add_text(slide, 2.65, ry + 0.85, 6.0, 0.25, drop_note,
                 color=color, size=10, bold=True, italic=True)
        add_text(slide, 2.65, ry + 1.05, 6.0, 0.20, kw_note,
                 color=MUTED, size=8, monospace=True)
        # Prediction box
        px, pw = 8.85, W - 0.60 - 8.55
        add_rect(slide, px, ry + 0.10, pw, row_h - 0.20,
                 fill=PANEL_FACE, line=color, line_w=1.4)
        add_text(slide, px + 0.10, ry + 0.15, pw - 0.20, 0.25,
                 "Prospective QA prediction:",
                 color=MUTED, size=9, bold=True)
        verdict = "✓  correct" if is_correct else "✗  incorrect"
        add_text(slide, px + 0.10, ry + 0.40, pw - 0.20, 0.30,
                 verdict, color=color, size=12, bold=True)
        add_text(slide, px + 0.10, ry + 0.72, pw - 0.20, 0.50,
                 pred, color=INK, size=10)


def build_good_gallery(prs: Presentation):
    W, H = 16.0, 10.0
    slide = _new_slide(prs, W, H)
    add_text(slide, 0.20, 0.10, W - 0.40, 0.40,
             "ForSeeBench gallery: forward-relevant evidence across target types",
             color=INK, size=22, bold=True, align="center")
    add_text(slide, 0.20, 0.55, W - 0.40, 0.35,
             "Each panel shows a retained ForSeeBench item: prior AD, hidden target, "
             "the four-option question, and the bridging keyword that ties prior context to the correct answer.",
             color=MUTED, size=11, italic=True, align="center")

    specs = [
        ("good_examples/example_03", "state_change",
         ["VOLDEMORT forces HARRY to bow."],
         "HARRY writhes in agony.",
         [("HARRY experiences intense pain.", True, "correct"),
          ("HARRY starts to laugh uncontrollably.", False, "unrelated"),
          ("HARRY suddenly gains superhuman strength.", False, "plausible_unsupported"),
          ("HARRY begins to float upwards.", False, "contradicts_context")],
         "forces … to bow"),
        ("good_examples/example_02", "participant_update",
         ["Mina moves towards Hansel.", "Mina kisses him."],
         "They embrace passionately.",
         [("Mina and Hansel kiss passionately.", True, "correct"),
          ("Hansel and Gretel kiss passionately.", False, "entity_swapped"),
          ("The witch appears and chases them away.", False, "plausible_unsupported"),
          ("The forest becomes covered in snow.", False, "contradicts_context")],
         "Mina · Hansel"),
        ("good_examples/example_04", "state_change",
         ["From a watch tower, Nantz and Imlay look over the extent of the destruction."],
         "Nantz peers through binoculars.",
         [("Nantz uses binoculars to observe the damage.", True, "correct"),
          ("Imlay takes over observing the damage.", False, "entity_swapped"),
          ("The destruction stops spreading.", False, "plausible_unsupported"),
          ("A new character appears on the scene.", False, "contradicts_context")],
         "watch tower · destruction"),
        ("good_examples/example_06", "participant_update",
         ["… which leads to a massive walk-in closet."],
         "Matty enters.",
         [("Matty walks into the walk-in closet.", True, "correct"),
          ("The walk-in closet suddenly disappears.", False, "unrelated"),
          ("Matty finds a secret room behind a bookshelf.", False, "plausible_unsupported"),
          ("A ghost appears in the walk-in closet.", False, "contradicts_context")],
         "walk-in closet"),
        ("good_examples/example_07", "spatial_consequence",
         ["As Stephen goes, Sarah smiles after him.",
          "Later, Sarah steps off the elevator on her floor."],
         "Crossing toward her bed.",
         [("Sarah walks towards her bed.", True, "correct"),
          ("Stephen enters the bedroom.", False, "already_happened"),
          ("Sarah sits on the couch.", False, "plausible_unsupported"),
          ("The doorbell rings.", False, "contradicts_context")],
         "Sarah · elevator → floor"),
        ("good_examples/example_10", "object_reveal",
         ["MOODY places the spider in front of HERMIONE.", "HERMIONE shakes her head."],
         "The spider dies.",
         [("Hermione observes the dead spider.", True, "correct"),
          ("The spider moves.", False, "already_happened"),
          ("Hermione talks to the spider.", False, "plausible_unsupported"),
          ("The spider is alive.", False, "contradicts_context")],
         "spider · Hermione"),
    ]
    cols, rows = 3, 2
    margin_x, margin_y = 0.30, 0.20
    grid_top, grid_bot = 1.05, 0.30
    panel_w = (W - margin_x * (cols + 1)) / cols
    panel_h = (H - grid_top - grid_bot - margin_y * (rows - 1)) / rows
    for idx, spec in enumerate(specs):
        r = idx // cols; c = idx % cols
        x = margin_x + c * (panel_w + margin_x)
        y = grid_top + r * (panel_h + margin_y)
        folder, ttype, prior, target, options, bridge = spec
        folder_p = FE / folder
        movie = short_movie(load_meta(folder_p)["movie"], limit=22)
        name_id = folder.split("/")[-1]
        add_panel(slide, x, y, panel_w, panel_h,
                  ec=GREEN, fc=PANEL_FACE,
                  header=f"{name_id}  ·  {movie}  ·  {ttype}",
                  header_color=GREEN, header_size=10,
                  header_height=0.32)

        # Frame
        fp = first_frame(folder_p)
        if fp:
            add_image(slide, fp, x + 0.15, y + 0.45, 1.45, 1.05)

        # Prior AD
        tx = x + 1.75
        tw = panel_w - 1.95
        add_text(slide, tx, y + 0.45, tw, 0.25,
                 "Prior AD context", color=MUTED, size=9, bold=True)
        yy_p = y + 0.72
        for line in prior:
            add_text(slide, tx, yy_p, 0.20, 0.28, "•",
                     color=INK, size=10, bold=True)
            rows_p = max(1, (len(line) + 49) // 50)
            add_text(slide, tx + 0.20, yy_p, tw - 0.20, 0.32 * rows_p,
                     line, color=INK, size=9.5)
            yy_p += 0.32 * rows_p + 0.05

        # Bridge chip
        add_rect(slide, tx, y + 1.55, tw, 0.40,
                 fill=PRIOR_HI, line=None)
        add_text(slide, tx + 0.10, y + 1.60, 0.80, 0.30,
                 "bridge:", color=ORANGE, size=9, bold=True)
        add_text(slide, tx + 0.95, y + 1.60, tw - 1.05, 0.30,
                 bridge, color=DARK_RED, size=10, bold=True, italic=True)

        # Hidden target
        add_text(slide, x + 0.15, y + 2.10, 1.60, 0.25,
                 "Hidden target AD:", color=DARK_RED, size=9, bold=True)
        add_text(slide, x + 1.75, y + 2.10, panel_w - 1.95, 0.45,
                 target, color=INK, size=9.5, italic=True)

        # Question + options
        add_text(slide, x + 0.15, y + 2.55, panel_w - 0.30, 0.30,
                 "Q: What happens next?", color=INK, size=10.5, bold=True)
        yy_o = y + 2.85
        for j, (opt, is_gold, tag) in enumerate(options):
            letter = chr(ord("A") + j)
            if is_gold:
                add_rect(slide, x + 0.10, yy_o - 0.03, panel_w - 0.20, 0.42,
                         fill=GOLD_HI, line=None)
            add_text(slide, x + 0.20, yy_o, 0.40, 0.32, f"({letter})",
                     color=GREEN if is_gold else INK, size=10, bold=True)
            add_text(slide, x + 0.65, yy_o, panel_w - 2.10, 0.42, opt,
                     color=INK, size=9.5)
            if is_gold:
                add_text(slide, x + panel_w - 1.30, yy_o, 0.30, 0.30, "✓",
                         color=GREEN, size=12, bold=True, align="right")
            add_text(slide, x + panel_w - 1.85, yy_o + 0.20, 1.70, 0.20,
                     tag, color=MUTED, size=7, monospace=True, align="right")
            yy_o += 0.42


def build_ad_source_grid(prs: Presentation):
    W, H = 16.0, 11.5
    slide = _new_slide(prs, W, H)
    add_text(slide, 0.20, 0.10, W - 0.40, 0.40,
             "AD-source substitution across three items",
             color=INK, size=22, bold=True, align="center")
    add_text(slide, 0.20, 0.55, W - 0.40, 0.40,
             "Same questions, options, and selected positions — only the AD source filling those positions changes. "
             "Human MAD-eval AD predicts correctly; NarrAD and AutoAD-Zero mispredict because the bridging entity or location is dropped.",
             color=MUTED, size=11, italic=True, align="center")

    rows = [
        ("baseline_comparison_examples/example_01",
         "Tracy groggily gets up and crosses the dim room.",
         "Tracy gets out of bed and walks across the room.",
         "In her room, Tracy sits up in bed and rubs her brow.",
         "Rebecca watches Sarah sleep, then wipes her face thoughtfully.",
         "Tracy covers her face with her hands in front of the mirror.",
         "Tracy gets out of bed and walks across the room.",
         "Tracy sits back down on the bed and continues to rub her brow.",
         "Tracy sits back down on the bed and continues to rub her brow.",
         "NarrAD swaps the protagonist; AutoAD-Zero changes the location."),
        ("baseline_comparison_examples/example_03",
         "Hansel then wraps a chain around her neck and hoists her up using the oven door as a pulley.",
         "Hansel uses the oven door as a pulley to lift the witch.",
         "Hansel heaves a hanging cauldron into the witches face knocking her down.",
         "The gingerbread house is in a tumultuous state.",
         "Hansel runs through the maze, looking around frantically.",
         "Hansel uses the oven door as a pulley to lift the witch.",
         "The witch transforms into a bird.",
         "The witch escapes by using the oven door as a ladder.",
         "Generated AD strips out both the actor (Hansel) and the object of action (witch / cauldron)."),
        ("baseline_comparison_examples/example_05",
         "And takes the receiver.",
         "MERRILL takes the receiver.",
         "MERRILL reaches out his hand. · MERRILL and GRAHAM both hold the receiver across the car roof.",
         "Merrill leans towards the rear, focusing intently on something.",
         "Graham reaches out and pulls Sara closer to him.",
         "MERRILL takes the receiver.",
         "GRAHAM takes the receiver.",
         "GRAHAM takes the receiver.",
         "Without a named co-actor (MERRILL holding the receiver), the answerer can't pick the right entity."),
    ]

    margin = 0.30
    top, bot = 1.10, 0.30
    panel_h = (H - top - bot - margin * (len(rows) - 1)) / len(rows)
    for i, spec in enumerate(rows):
        (folder, target, correct, h_ad, n_ad, a_ad,
         h_pred, n_pred, a_pred, drop_summary) = spec
        movie = short_movie(load_meta(FE / folder)["movie"], limit=28)
        y = top + i * (panel_h + margin)
        add_panel(slide, 0.30, y, W - 0.60, panel_h,
                  ec=BLUE, fc=PANEL_FACE,
                  header=f"{folder.split('/')[-1]}  ·  {movie}",
                  header_color=BLUE, header_size=11,
                  header_height=0.32)

        # Q + hidden target
        add_text(slide, 0.50, y + 0.45, 9, 0.30,
                 "Q:  What happens next?", color=INK, size=11, bold=True)
        add_text(slide, 0.50, y + 0.75, 1.65, 0.25,
                 "Hidden target AD:", color=DARK_RED, size=9, bold=True)
        add_text(slide, 2.20, y + 0.75, W - 3.0, 0.40,
                 target, color=INK, size=9.5, italic=True)

        # Correct option banner
        add_rect(slide, 0.40, y + 1.20, W - 0.80, 0.42,
                 fill=GOLD_HI, line=None)
        add_text(slide, 0.55, y + 1.25, 1.80, 0.30,
                 "✓ correct option:", color=GREEN, size=10, bold=True)
        add_text(slide, 2.30, y + 1.25, W - 3.0, 0.32,
                 correct, color=INK, size=10.5, bold=True)

        # Drop summary chip
        add_rect(slide, 0.40, y + 1.68, W - 0.80, 0.36,
                 fill=PRIOR_HI, line=None)
        add_text(slide, 0.55, y + 1.72, 1.65, 0.28,
                 "evidence gap:", color=ORANGE, size=9.5, bold=True)
        add_text(slide, 2.25, y + 1.72, W - 3.0, 0.30,
                 drop_summary, color=DARK_RED, size=10, italic=True)

        # Three columns
        col_w = (W - 0.80) / 3 - 0.10
        sources = [
            ("Human MAD-eval AD", h_ad, h_pred, h_pred == correct, GREEN, GOLD_HI),
            ("NarrAD",            n_ad, n_pred, n_pred == correct, RED, RED_HI),
            ("AutoAD-Zero",       a_ad, a_pred, a_pred == correct, RED, RED_HI),
        ]
        col_y = y + 2.15
        for j, (label, ad_text, pred, is_correct, color, hi_color) in enumerate(sources):
            x = 0.40 + j * (col_w + 0.10)
            # Header strip
            add_rect(slide, x, col_y, col_w, 0.32,
                     fill=color, line=None)
            add_text(slide, x, col_y, col_w, 0.32, label,
                     color=WHITE, size=10, bold=True, align="center", vertical="middle")
            # AD text box
            add_rect(slide, x, col_y + 0.32, col_w, 0.65,
                     fill=PANEL_FACE, line=color, line_w=0.7)
            add_text(slide, x + 0.10, col_y + 0.36, col_w - 0.20, 0.20,
                     "AD filling selected position:",
                     color=MUTED, size=8, bold=True)
            add_text(slide, x + 0.10, col_y + 0.55, col_w - 0.20, 0.40,
                     ad_text, color=INK, size=9)
            # Prediction box
            add_rect(slide, x, col_y + 0.97, col_w, 0.55,
                     fill=hi_color, line=color, line_w=1.0)
            add_text(slide, x + 0.10, col_y + 1.02, col_w - 0.20, 0.20,
                     "prospective QA prediction:",
                     color=MUTED, size=8, bold=True)
            verdict = "✓ correct" if is_correct else "✗ incorrect"
            add_text(slide, x + 0.10, col_y + 1.18, col_w - 0.20, 0.20,
                     verdict, color=color, size=10, bold=True)
            add_text(slide, x + 0.10, col_y + 1.34, col_w - 0.20, 0.18,
                     pred, color=INK, size=9)


def build_target_types(prs: Presentation):
    W, H = 16.0, 9.5
    slide = _new_slide(prs, W, H)
    add_text(slide, 0.20, 0.10, W - 0.40, 0.40,
             "Target-type taxonomy: ForSeeBench is broader than action prediction",
             color=INK, size=22, bold=True, align="center")
    add_text(slide, 0.20, 0.55, W - 0.40, 0.30,
             "Five of the six retained target types, one example each. "
             "(visible_text_update, n=7, omitted for compact layout.)",
             color=MUTED, size=11, italic=True, align="center")

    type_specs = [
        ("good_examples/example_02", "participant_update", BLUE, BLUE_HI,
         ["Mina moves towards Hansel.", "Mina kisses him."],
         "They embrace passionately.",
         "Mina and Hansel kiss passionately.",
         "A new participant enters or becomes the focus of the scene."),
        ("good_examples/example_09", "action_transition", PURPLE, PURPLE_HI,
         ["GRAHAM pushes his hand against the door lit by MORGAN torch."],
         "GRAHAM then turns to face the others.",
         "GRAHAM prepares to enter the room.",
         "A character transitions from one action to the next."),
        ("good_examples/example_03", "state_change", GREEN, GOLD_HI,
         ["VOLDEMORT forces HARRY to bow."],
         "HARRY writhes in agony.",
         "HARRY experiences intense pain.",
         "A visible state of a character or object changes."),
        ("good_examples/example_07", "spatial_consequence", TEAL, TEAL_HI,
         ["As Stephen goes, Sarah smiles after him.",
          "Later, Sarah steps off the elevator on her floor."],
         "Crossing toward her bed.",
         "Sarah walks towards her bed.",
         "A spatial movement / location change implied by the prior AD."),
        ("good_examples/example_10", "object_reveal", ORANGE, ORANGE_HI,
         ["MOODY places the spider in front of HERMIONE.", "HERMIONE shakes her head."],
         "The spider dies.",
         "Hermione observes the dead spider.",
         "An object becomes narratively important / changes status."),
    ]
    cols, rows = 3, 2
    margin_x, margin_y = 0.30, 0.25
    grid_top, grid_bot = 1.05, 0.30
    panel_w = (W - margin_x * (cols + 1)) / cols
    panel_h = (H - grid_top - grid_bot - margin_y * (rows - 1)) / rows

    for idx, spec in enumerate(type_specs):
        r = idx // cols; c = idx % cols
        x = margin_x + c * (panel_w + margin_x)
        y = grid_top + r * (panel_h + margin_y)
        folder, ttype, color, hi_color, prior, target, correct, interp = spec
        folder_p = FE / folder
        movie = short_movie(load_meta(folder_p)["movie"], limit=30)
        add_panel(slide, x, y, panel_w, panel_h,
                  ec=color, fc=PANEL_FACE,
                  header=f"target type:  {ttype}",
                  header_color=color, header_size=12,
                  header_height=0.40)
        add_text(slide, x + 0.20, y + 0.50, panel_w - 0.30, 0.25,
                 f"{folder.split('/')[-1]} · {movie}",
                 color=MUTED, size=9, italic=True)
        # Frame
        fp = first_frame(folder_p)
        if fp:
            add_image(slide, fp, x + 0.20, y + 0.85, 1.70, 1.20)
        # Prior box
        tx = x + 2.05
        tw = panel_w - 2.25
        add_text(slide, tx, y + 0.85, tw, 0.25,
                 "Prior AD context", color=MUTED, size=10, bold=True)
        add_rect(slide, tx, y + 1.10, tw, 1.05,
                 fill=hi_color, line=None)
        prior_text = "\n".join("• " + l for l in prior)
        add_text(slide, tx + 0.10, y + 1.15, tw - 0.20, 0.95,
                 prior_text, color=INK, size=10, line_spacing=1.20)
        # Hidden target
        add_text(slide, x + 0.20, y + 2.30, panel_w - 0.30, 0.25,
                 "Hidden target AD:", color=DARK_RED, size=9, bold=True)
        add_text(slide, x + 0.20, y + 2.55, panel_w - 0.30, 0.45,
                 f"“{target}”", color=INK, size=10, italic=True)
        # Correct banner
        add_rect(slide, x + 0.20, y + 3.00, panel_w - 0.40, 0.65,
                 fill=hi_color, line=color, line_w=1.4)
        add_text(slide, x + 0.30, y + 3.05, panel_w - 0.60, 0.25,
                 "✓ correct option:",
                 color=color, size=10, bold=True)
        add_text(slide, x + 0.30, y + 3.30, panel_w - 0.60, 0.30,
                 correct, color=INK, size=10.5, bold=True)
        # Interpretation
        add_text(slide, x + 0.20, y + 3.75, panel_w - 0.30, 0.30,
                 interp, color=MUTED, size=9, italic=True)

    # Stats panel (sixth cell)
    idx = len(type_specs)
    r = idx // cols; c = idx % cols
    x = margin_x + c * (panel_w + margin_x)
    y = grid_top + r * (panel_h + margin_y)
    add_panel(slide, x, y, panel_w, panel_h,
              ec=MUTED, fc=PANEL_TINT,
              header="benchmark composition (787 items)",
              header_color=MUTED, header_size=11, header_height=0.40)
    add_text(slide, x + 0.30, y + 0.55, panel_w * 0.55, 0.25,
             "Target type", color=INK, size=10.5, bold=True)
    add_text(slide, x + panel_w - 1.30, y + 0.55, 1.0, 0.25,
             "Count", color=INK, size=10.5, bold=True, align="right")
    counts = [
        ("participant_update",  "379"),
        ("action_transition",   "196"),
        ("state_change",        "147"),
        ("object_reveal",       " 30"),
        ("spatial_consequence", " 28"),
        ("visible_text_update", "  7"),
    ]
    yy = y + 0.85
    for tt, n in counts:
        add_text(slide, x + 0.30, yy, panel_w * 0.55, 0.25,
                 tt, color=INK, size=10, monospace=True)
        add_text(slide, x + panel_w - 1.30, yy, 1.0, 0.25,
                 n, color=INK, size=10, monospace=True, align="right")
        yy += 0.32

    add_text(slide, x + 0.30, yy + 0.30, panel_w - 0.50, 0.25,
             "Adaptive context length",
             color=INK, size=10, bold=True)
    add_text(slide, x + 0.30, yy + 0.55, panel_w - 0.50, 0.25,
             "median = 1 clip · mean = 2.01 · max = 9",
             color=MUTED, size=9, italic=True)


def build_evidence_chain(prs: Presentation):
    W, H = 13.5, 11.0
    slide = _new_slide(prs, W, H)
    add_text(slide, 0.20, 0.10, W - 0.40, 0.40,
             "Forward-relevant evidence: a single keyword bridges prior AD to the correct option",
             color=INK, size=20, bold=True, align="center")
    add_text(slide, 0.20, 0.55, W - 0.40, 0.40,
             "Each row reads left-to-right: prior AD context  →  bridging keyword  →  correct option (paraphrases the hidden target). "
             "The bridge keyword is a single object, participant, or location.",
             color=MUTED, size=11, italic=True, align="center")

    chain_specs = [
        ("good_examples/example_03",
         ["VOLDEMORT forces HARRY to bow."],
         "forces … to bow",
         "HARRY writhes in agony.",
         "HARRY experiences intense pain.",
         ORANGE, GREEN, GOLD_HI),
        ("good_examples/example_06",
         ["… which leads to a massive walk-in closet."],
         "walk-in closet",
         "Matty enters.",
         "Matty walks into the walk-in closet.",
         ORANGE, BLUE, BLUE_HI),
        ("good_examples/example_04",
         ["From a watch tower, Nantz and Imlay look over the extent of the destruction."],
         "watch tower · destruction",
         "Nantz peers through binoculars.",
         "Nantz uses binoculars to observe the damage.",
         ORANGE, TEAL, TEAL_HI),
        ("good_examples/example_10",
         ["MOODY places the spider in front of HERMIONE.", "HERMIONE shakes her head."],
         "spider · Hermione",
         "The spider dies.",
         "Hermione observes the dead spider.",
         ORANGE, PURPLE, PURPLE_HI),
    ]
    margin = 0.20
    top, bot = 1.05, 0.30
    panel_h = (H - top - bot - margin * (len(chain_specs) - 1)) / len(chain_specs)

    for i, spec in enumerate(chain_specs):
        (folder, prior_lines, bridge, target, correct,
         arrow_color, panel_color, panel_hi) = spec
        movie = short_movie(load_meta(FE / folder)["movie"], limit=30)
        y = top + i * (panel_h + margin)
        add_panel(slide, 0.30, y, W - 0.60, panel_h,
                  ec=panel_color, fc=PANEL_FACE,
                  header=f"{folder.split('/')[-1]}  ·  {movie}",
                  header_color=panel_color, header_size=11,
                  header_height=0.32)
        # Frame thumbnail in corner
        fp = first_frame(FE / folder)
        if fp:
            add_image(slide, fp, 0.45, y + 0.45, 1.40, 0.55)
        add_text(slide, 1.95, y + 0.55, 8.0, 0.30,
                 "context frame (visualization only — answerer is text-only)",
                 color=MUTED, size=9, italic=True)

        # Pill 1: Prior AD
        pill_y = y + 1.20
        pill_h = panel_h - 1.50
        add_rect(slide, 0.45, pill_y, 3.40, pill_h,
                 fill=panel_hi, line=panel_color, line_w=1.0)
        add_text(slide, 0.55, pill_y + 0.12, 3.20, 0.25,
                 "Prior AD context",
                 color=panel_color, size=10, bold=True)
        prior_text = "\n".join("• " + l for l in prior_lines)
        add_text(slide, 0.55, pill_y + 0.40, 3.20, pill_h - 0.55,
                 prior_text, color=INK, size=10, line_spacing=1.25)

        # Pill 2: Bridge
        bridge_x = 4.30
        bridge_w = 2.50
        add_rect(slide, bridge_x, pill_y + 0.20,
                 bridge_w, pill_h - 0.40,
                 fill=PRIOR_HI, line=arrow_color, line_w=1.6)
        add_text(slide, bridge_x, pill_y + 0.30, bridge_w, 0.30,
                 "bridge keyword",
                 color=arrow_color, size=10, bold=True, align="center")
        add_text(slide, bridge_x, pill_y + (pill_h - 0.40) / 2 - 0.05,
                 bridge_w, 0.50,
                 bridge, color=DARK_RED, size=12, bold=True, italic=True,
                 align="center", vertical="middle")

        # Pill 3: Correct
        ans_x = 7.20
        ans_w = W - ans_x - 0.45
        add_rect(slide, ans_x, pill_y, ans_w, pill_h,
                 fill=GOLD_HI, line=panel_color, line_w=1.6)
        add_text(slide, ans_x + 0.10, pill_y + 0.12, ans_w - 0.20, 0.25,
                 "✓  correct option (paraphrases hidden target)",
                 color=panel_color, size=10, bold=True)
        add_text(slide, ans_x + 0.10, pill_y + 0.40, ans_w - 0.20, pill_h - 0.55,
                 correct, color=INK, size=11, bold=True)
        # Hidden target line below
        add_text(slide, ans_x + 0.10, y + panel_h - 0.50, ans_w - 0.20, 0.20,
                 "hidden target AD:", color=DARK_RED, size=8.5, bold=True)
        add_text(slide, ans_x + 0.10, y + panel_h - 0.30, ans_w - 0.20, 0.25,
                 f"“{target}”", color=INK, size=9.5, italic=True)

        # Arrows between pills
        mid_y = pill_y + pill_h / 2
        add_arrow(slide, 3.85, mid_y, 4.30, mid_y, color=arrow_color, width=2.4)
        add_arrow(slide, 6.80, mid_y, 7.20, mid_y, color=arrow_color, width=2.4)


# ---------------------------------------------------------------------------
# Driver: write each figure as its own deck + a combined deck
# ---------------------------------------------------------------------------
BUILDERS = [
    ("fig_teaser.pptx",         build_teaser),
    ("fig_filter.pptx",         build_filter),
    ("fig_ad_source.pptx",      build_ad_source),
    ("fig_good_gallery.pptx",   build_good_gallery),
    ("fig_ad_source_grid.pptx", build_ad_source_grid),
    ("fig_target_types.pptx",   build_target_types),
    ("fig_evidence_chain.pptx", build_evidence_chain),
]


def main() -> None:
    print("Building editable PPTX figures →", OUT)

    # Single-figure decks
    for name, builder in BUILDERS:
        prs = Presentation()
        builder(prs)
        prs.save(OUT / name)
        print(f"  wrote {OUT / name}")

    # Combined deck — one slide per figure, identical content
    combo = Presentation()
    for _, builder in BUILDERS:
        builder(combo)
    combo.save(OUT / "all_figures.pptx")
    print(f"  wrote {OUT / 'all_figures.pptx'}")
    print("Done.")


if __name__ == "__main__":
    main()
