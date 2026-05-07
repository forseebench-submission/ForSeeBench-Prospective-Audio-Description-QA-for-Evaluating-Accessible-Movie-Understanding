"""Build the ForSeeBench construction-pipeline figure (clean version).

Outputs:
  figures/paper_figures/fig_construction_pipeline.{pdf,png}
  figures/paper_figures_pptx/fig_construction_pipeline.pptx

The figure is a single horizontal flow with five stages, downward rejection
branches, and a green final-output card. All counts are taken from
Table~1 of the paper (data/processed/all_movies/eval_all10.jsonl) so the
figure stays in sync with the released benchmark.
"""
from __future__ import annotations

from pathlib import Path

import matplotlib.patches as patches
import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch

ROOT = Path("/jumbo/jinlab/Sally/ForSeeBench")
OUT_PNG = ROOT / "figures/paper_figures"
OUT_PPTX_DIR = ROOT / "figures/paper_figures_pptx"
OUT_PNG.mkdir(parents=True, exist_ok=True)
OUT_PPTX_DIR.mkdir(parents=True, exist_ok=True)

# Palette (matches the rest of the paper figures)
INK = "#1F2933"
MUTED = "#6B7280"
SUBTLE = "#9CA3AF"
LIGHT = "#F5F5F7"
BLUE = "#3D6FB6"
DARK_BLUE = "#1E3A8A"
BLUE_HI = "#E8F0FB"
GREEN = "#3F8F5C"
DARK_GREEN = "#1E5631"
GREEN_HI = "#DDEFD2"
RED = "#C0413E"
DARK_RED = "#7B1F1C"
RED_HI = "#FBE3E1"
ORANGE = "#D97A2D"


# ---------------------------------------------------------------------------
# Stage data — every count is auditable against the paper / pipeline files.
# ---------------------------------------------------------------------------
STAGES = [
    {
        "n": 1,
        "title": "Source AD",
        "desc": "Ordered AD clips\nfrom 10 MAD-eval\nmovies.",
        "count_label": "AD clips",
        "count_value": "6,520",
        "rejection": None,
    },
    {
        "n": 2,
        "title": "Search windows",
        "desc": "Sliding 10-clip\nblocks (one per\nstarting AD clip).",
        "count_label": "search blocks",
        "count_value": "6,520",
        "rejection": None,
    },
    {
        "n": 3,
        "title": "Target + context",
        "desc": "Qwen2.5-VL picks a\nhidden future target\nand its strictly-prior\nevidence clips.",
        "count_label": "candidates",
        "count_value": "799",
        "rejection": {
            "count": "332 blocks",
            "reasons": "no valid target /\ninsufficient buildup",
        },
    },
    {
        "n": 4,
        "title": "Question + 4 options",
        "desc": "Qwen2.5-VL writes\none MCQ:\n1 correct + 3 typed\ndistractors.",
        "count_label": "candidates",
        "count_value": "792",
        "rejection": {
            "count": "7 blocks",
            "reasons": "schema-invalid /\nunanswerable target",
        },
    },
    {
        "n": 5,
        "title": "Validation gates",
        "desc": "Qwen2.5-VL scores\nconfidence,\nevidence sufficiency,\ndistractor quality\n(all ≥ 0.7).",
        "count_label": "kept",
        "count_value": "787",
        "rejection": {
            "count": "5 items",
            "reasons": "distractor quality\n= 0.0",
        },
    },
]

TARGET_TYPES = [
    ("participant_update", 379),
    ("action_transition",  196),
    ("state_change",       147),
    ("object_reveal",       30),
    ("spatial_consequence", 28),
    ("visible_text_update",  7),
]

DISTRACTOR_TYPES = [
    ("already_happened",       578),
    ("plausible_unsupported",  779),
    ("contradicts_context",    741),
    ("entity_swapped",         184),
    ("unrelated",               79),
]


# ---------------------------------------------------------------------------
# Matplotlib version
# ---------------------------------------------------------------------------
def _round_box(ax, x, y, w, h, *, ec, fc, lw=1.2, r=0.020):
    ax.add_patch(FancyBboxPatch(
        (x, y), w, h,
        boxstyle=f"round,pad=0.0,rounding_size={r}",
        lw=lw, ec=ec, fc=fc, zorder=2,
    ))


def _arrow(ax, x1, y1, x2, y2, *, color=INK, lw=1.4, ls="-"):
    ax.add_patch(FancyArrowPatch(
        (x1, y1), (x2, y2),
        arrowstyle="-|>", mutation_scale=12,
        linewidth=lw, color=color, linestyle=ls, zorder=4,
    ))


def make_matplotlib() -> None:
    plt.rcParams.update({
        "font.family": "DejaVu Sans",
        "font.size": 9,
        "pdf.fonttype": 42, "ps.fonttype": 42,
    })

    # Wider canvas + more vertical room. The funnel is the only content.
    fig = plt.figure(figsize=(9.5, 5.4))
    fig.patch.set_facecolor("white")
    ax = fig.add_axes([0, 0, 1, 1]); ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    ax.axis("off")

    # ---- Title ----
    ax.text(0.50, 0.965,
            "ForSeeBench construction pipeline",
            ha="center", va="center", fontsize=13, fontweight="bold", color=INK)
    ax.text(0.50, 0.928,
            "10 MAD-eval movies  →  787 prospective QA items "
            "(every rejection is logged for audit).",
            ha="center", va="center", fontsize=8.6, color=MUTED, style="italic")

    # ---- Stage geometry ----
    n = len(STAGES)
    margin = 0.018
    stage_y_top = 0.880
    stage_y_bot = 0.555
    stage_h = stage_y_top - stage_y_bot
    avail = 1.0 - 2 * margin
    gap = 0.018
    stage_w = (avail - gap * (n - 1)) / n
    stage_x = [margin + i * (stage_w + gap) for i in range(n)]
    arrow_y = (stage_y_top + stage_y_bot) / 2

    # Arrows between stages
    for i in range(n - 1):
        _arrow(ax, stage_x[i] + stage_w, arrow_y,
               stage_x[i + 1], arrow_y, color=BLUE, lw=1.4)

    # Stages
    for i, st in enumerate(STAGES):
        x = stage_x[i]
        # Body
        _round_box(ax, x, stage_y_bot, stage_w, stage_h,
                   ec=BLUE, fc="#FFFFFF", lw=1.4, r=0.020)
        # Stage number badge (top-left)
        badge_w = 0.028
        badge_h = 0.040
        _round_box(ax, x + 0.010, stage_y_top - badge_h - 0.010,
                   badge_w, badge_h,
                   ec=BLUE, fc=BLUE, lw=0, r=0.014)
        ax.text(x + 0.010 + badge_w / 2,
                stage_y_top - badge_h / 2 - 0.010,
                str(st["n"]),
                ha="center", va="center", fontsize=10,
                color="white", fontweight="bold", zorder=3)
        # Stage title (right of badge, single line)
        title = st["title"].replace("\n", " ")
        ax.text(x + 0.045, stage_y_top - 0.030,
                title,
                fontsize=8.6, fontweight="bold", color=DARK_BLUE, va="center")
        # Description
        ax.text(x + 0.012, stage_y_top - 0.080, st["desc"],
                fontsize=7.4, color=INK, va="top", linespacing=1.32)

        # Count
        ax.text(x + stage_w / 2, stage_y_bot + 0.060,
                st["count_value"],
                ha="center", va="center", fontsize=15,
                fontweight="bold", color=DARK_BLUE)
        ax.text(x + stage_w / 2, stage_y_bot + 0.022,
                st["count_label"],
                ha="center", va="center", fontsize=7.2,
                color=MUTED, style="italic")

        # Rejection branch
        if st.get("rejection"):
            # Down arrow
            _arrow(ax, x + stage_w / 2, stage_y_bot - 0.008,
                   x + stage_w / 2, 0.470, color=RED, lw=1.0,
                   ls=(0, (3, 2)))
            # Red box
            rej_x = x + 0.010
            rej_w = stage_w - 0.020
            rej_y_top = 0.470
            rej_y_bot = 0.330
            _round_box(ax, rej_x, rej_y_bot, rej_w, rej_y_top - rej_y_bot,
                       ec=RED, fc=RED_HI, lw=0.9, r=0.014)
            ax.text(rej_x + rej_w / 2, rej_y_top - 0.030,
                    f"rejected: {st['rejection']['count']}",
                    ha="center", va="center", fontsize=8.0,
                    color=DARK_RED, fontweight="bold")
            ax.text(rej_x + rej_w / 2, rej_y_top - 0.080,
                    st["rejection"]["reasons"],
                    ha="center", va="top", fontsize=7.0,
                    color=DARK_RED, style="italic", linespacing=1.30)

    # ---- Final output card ----
    out_y_top = 0.265
    out_y_bot = 0.060
    out_h = out_y_top - out_y_bot
    out_x = 0.020
    out_w = 1 - 2 * out_x

    # Vertical green arrow from stage 5 into the output card
    last_x = stage_x[-1] + stage_w / 2
    _arrow(ax, last_x, stage_y_bot - 0.008,
           last_x, out_y_top + 0.005, color=GREEN, lw=2.0)

    _round_box(ax, out_x, out_y_bot, out_w, out_h,
               ec=GREEN, fc=GREEN_HI, lw=1.6, r=0.020)
    ax.text(out_x + 0.110, out_y_bot + out_h * 0.62,
            "787",
            ha="center", va="center", fontsize=32,
            fontweight="bold", color=DARK_GREEN)
    ax.text(out_x + 0.110, out_y_bot + out_h * 0.20,
            "released\nbenchmark items",
            ha="center", va="center", fontsize=8.0, color=DARK_GREEN,
            style="italic", linespacing=1.30)

    info_x = out_x + 0.230
    ax.text(info_x, out_y_top - 0.030,
            "Each item stores:",
            fontsize=9.4, fontweight="bold", color=DARK_GREEN, va="top")
    ax.text(info_x, out_y_top - 0.072,
            "•  prior-AD context ids   •  hidden target AD\n"
            "•  question + 4 options    •  typed distractor labels (1 of 5)\n"
            "•  evidence spans matching strictly prior AD text\n"
            "•  target type (1 of 6), reasoning type, continuity type\n"
            "•  Qwen confidence / evidence sufficiency / distractor quality (≥ 0.7)",
            fontsize=7.6, color=INK, va="top", linespacing=1.45)

    pdf = OUT_PNG / "fig_construction_pipeline.pdf"
    png = OUT_PNG / "fig_construction_pipeline.png"
    fig.savefig(pdf, bbox_inches="tight")
    fig.savefig(png, dpi=300, bbox_inches="tight")
    plt.close(fig)
    print(f"  wrote {pdf}  +  {png}")


# ---------------------------------------------------------------------------
# PPTX version
# ---------------------------------------------------------------------------
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_rect(slide, x, y, w, h, *, fill, line, line_w=1.0, rounded=True,
              corner=0.10):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    if rounded:
        s.adjustments[0] = corner
    if fill is not None:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def _add_text(slide, x, y, w, h, text, *, color, size, bold=False,
              italic=False, align="left", vertical="middle", monospace=False,
              line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if vertical == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif vertical == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = "Consolas" if monospace else "DejaVu Sans"
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return box


def _add_arrow(slide, x1, y1, x2, y2, *, color, width=1.5, dashed=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = line.line._get_or_add_ln()
    tail_end = ln.find(qn("a:tailEnd"))
    if tail_end is None:
        tail_end = etree.SubElement(ln, qn("a:tailEnd"))
    tail_end.set("type", "triangle")
    tail_end.set("w", "med")
    tail_end.set("len", "med")
    if dashed:
        prst_dash = ln.find(qn("a:prstDash"))
        if prst_dash is None:
            prst_dash = etree.SubElement(ln, qn("a:prstDash"))
        prst_dash.set("val", "dash")
    return line


def make_pptx() -> None:
    prs = Presentation()
    W, H = 13.5, 7.0
    prs.slide_width = Inches(W); prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    INK_ = _rgb(INK); MUTED_ = _rgb(MUTED); SUBTLE_ = _rgb(SUBTLE)
    BLUE_ = _rgb(BLUE); BLUE_HI_ = _rgb(BLUE_HI); DARK_BLUE_ = _rgb(DARK_BLUE)
    RED_ = _rgb(RED); RED_HI_ = _rgb(RED_HI); DARK_RED_ = _rgb(DARK_RED)
    GREEN_ = _rgb(GREEN); GREEN_HI_ = _rgb(GREEN_HI); DARK_GREEN_ = _rgb(DARK_GREEN)
    WHITE = _rgb("#FFFFFF")

    # Title
    _add_text(slide, 0.30, 0.20, W - 0.60, 0.40,
              "ForSeeBench construction pipeline",
              color=INK_, size=22, bold=True, align="center")
    _add_text(slide, 0.30, 0.62, W - 0.60, 0.30,
              "10 MAD-eval movies  →  787 prospective QA items, "
              "with each rejection logged for audit.",
              color=MUTED_, size=11, italic=True, align="center")

    # Stage layout
    margin_x = 0.30
    n = len(STAGES)
    gap = 0.20
    avail = W - 2 * margin_x
    stage_w = (avail - gap * (n - 1)) / n
    stage_y = 1.20
    stage_h = 2.20
    stage_xs = [margin_x + i * (stage_w + gap) for i in range(n)]

    # Arrows between stages
    for i in range(n - 1):
        _add_arrow(slide,
                   stage_xs[i] + stage_w, stage_y + stage_h / 2,
                   stage_xs[i + 1], stage_y + stage_h / 2,
                   color=BLUE_, width=2.4)

    # Stages
    for i, st in enumerate(STAGES):
        x = stage_xs[i]
        # Body
        _add_rect(slide, x, stage_y, stage_w, stage_h,
                  fill=WHITE, line=BLUE_, line_w=1.6, corner=0.08)
        # Ribbon header
        _add_rect(slide, x, stage_y, stage_w, 0.40,
                  fill=BLUE_, line=None, corner=0.18)
        _add_text(slide, x + 0.12, stage_y, 1.0, 0.40,
                  f"Stage {st['n']}",
                  color=WHITE, size=11, bold=True, vertical="middle")
        _add_text(slide, x + 1.10, stage_y, stage_w - 1.20, 0.40,
                  st["title"].replace("\n", " · "),
                  color=WHITE, size=11, bold=True, align="right",
                  vertical="middle")

        # Description
        _add_text(slide, x + 0.15, stage_y + 0.45, stage_w - 0.30, 1.10,
                  st["desc"], color=INK_, size=10, vertical="top",
                  line_spacing=1.25)

        # Count
        _add_text(slide, x, stage_y + stage_h - 0.85, stage_w, 0.45,
                  st["count_value"],
                  color=DARK_BLUE_, size=22, bold=True, align="center",
                  vertical="middle")
        _add_text(slide, x, stage_y + stage_h - 0.40, stage_w, 0.30,
                  st["count_label"],
                  color=MUTED_, size=9, italic=True, align="center",
                  vertical="middle")

        # Rejection branch
        if st.get("rejection"):
            # Down arrow
            _add_arrow(slide,
                       x + stage_w / 2, stage_y + stage_h + 0.05,
                       x + stage_w / 2, stage_y + stage_h + 0.55,
                       color=RED_, width=1.2, dashed=True)
            rej_y = stage_y + stage_h + 0.60
            rej_h = 1.05
            rej_w = stage_w - 0.10
            rej_x = x + (stage_w - rej_w) / 2
            _add_rect(slide, rej_x, rej_y, rej_w, rej_h,
                      fill=RED_HI_, line=RED_, line_w=1.0, corner=0.12)
            _add_text(slide, rej_x, rej_y + 0.05, rej_w, 0.30,
                      f"rejected: {st['rejection']['count']}",
                      color=DARK_RED_, size=10, bold=True, align="center",
                      vertical="middle")
            _add_text(slide, rej_x + 0.10, rej_y + 0.40, rej_w - 0.20, rej_h - 0.50,
                      st["rejection"]["reasons"],
                      color=DARK_RED_, size=9, italic=True, align="center",
                      vertical="top", line_spacing=1.25)

    # ---- Final output card ----
    out_y = stage_y + stage_h + 1.85
    out_h = 1.30
    out_x = margin_x
    out_w = W - 2 * margin_x

    # Down-arrow from stage 5 to output card
    last_x = stage_xs[-1] + stage_w / 2
    _add_arrow(slide,
               last_x, stage_y + stage_h + 0.05,
               last_x, out_y - 0.05,
               color=GREEN_, width=2.6)

    _add_rect(slide, out_x, out_y, out_w, out_h,
              fill=GREEN_HI_, line=GREEN_, line_w=1.6, corner=0.08)
    _add_text(slide, out_x, out_y + 0.10, 2.50, 0.65,
              "787",
              color=DARK_GREEN_, size=44, bold=True, align="center",
              vertical="middle")
    _add_text(slide, out_x, out_y + 0.85, 2.50, 0.35,
              "released benchmark items",
              color=DARK_GREEN_, size=10, italic=True, align="center",
              vertical="middle")
    _add_text(slide, out_x + 2.70, out_y + 0.10, 1.80, 0.30,
              "Each item stores:",
              color=DARK_GREEN_, size=11, bold=True, vertical="top")
    _add_text(slide, out_x + 2.70, out_y + 0.40, out_w - 2.85, out_h - 0.45,
              "•  prior-AD context ids   •  hidden target AD\n"
              "•  question + 4 options   •  typed distractor labels\n"
              "•  evidence spans (must match strictly prior AD text)\n"
              "•  target type, reasoning type, continuity type\n"
              "•  Qwen confidence / evidence sufficiency / distractor quality",
              color=INK_, size=9.5, vertical="top", line_spacing=1.30)

    prs.save(OUT_PPTX_DIR / "fig_construction_pipeline.pptx")
    print(f"  wrote {OUT_PPTX_DIR / 'fig_construction_pipeline.pptx'}")


def main() -> None:
    print("Building construction-pipeline figure…")
    make_matplotlib()
    make_pptx()
    print("Done.")


if __name__ == "__main__":
    main()
